import sys
import tkinter as tk
import convertapi
from tkinter import ttk, filedialog, messagebox
from tkinterdnd2 import TkinterDnD, DND_FILES
import json
import requests
import time
import mimetypes
import os
import re
import subprocess
import logging
from pathlib import Path
import threading
import queue
import time
import openai
import deepl
import pptx
import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR, MSO_COLOR_TYPE
from PIL import Image
from pydub import AudioSegment
from core.tool_descriptions import get_short_description, get_tool_info, get_quick_tips
from core.transcription import AudioTranscriptionCore

# Run migration if needed (before other initializations)
try:
    subprocess.run([sys.executable, "migrate_secret.py", "--auto"], check=False)
except Exception:
    pass  # Migration script might not exist or fail, continue anyway


# --- Constants ---
SUPPORTED_LANGUAGES_FILE = "supported_languages.json"
API_KEYS_FILE = "api_keys.json"

# --- Setup Logging ---
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")

# Import consolidated core modules
from core.config import ConfigManager
from core.services import ServiceManager, create_service_manager
from core.processors import (
    create_translation_processor, create_audio_processor, create_conversion_processor,
    ProgressReporter, ProcessorConfig, ProcessingStatus
)
from core.task_manager import TaskManager, QueueProgressAdapter
from core.validation import validate_file_size, validate_language_code
from core.file_utils import should_skip_processing, get_file_size_mb

class ToolBase:
    """Base class for all tools in the application."""
    
    def __init__(self, master, config_manager, progress_queue):
        """Initialize the tool base class."""
        self.master = master
        self.config_manager = config_manager
        self.progress_queue = progress_queue
        self.input_paths = []
        self.output_path = None
        self.supported_languages = self.config_manager.get_languages()
        
        # Initialize service manager for consolidated API access
        self.service_manager = create_service_manager(config_manager)
        
        # Create progress reporter for GUI
        self.progress_reporter = ProgressReporter(callback=self.send_progress_update)
        
        # Add selection mode variable
        self.selection_mode = tk.StringVar(value="file")  # "file" or "folder"
        
        # Add output checking option
        self.check_output_exists = tk.BooleanVar(value=True)  # Default to checking if output exists
        
        # Define supported extensions for the tool (to be overridden by child classes)
        self.supported_extensions = set()
        
        # Initialize display attributes
        self.input_paths_display = None
        self.output_path_display = None

        self.stop_flag = threading.Event()
        self.processing_thread = None

    def stop_processing(self):
        """Signals the processing to stop."""
        if self.stop_flag:
            self.stop_flag.set()
            self.send_progress_update("Stopping processing... Please wait.")
            logging.info("Stop processing requested")

    def create_selection_mode_controls(self, parent_frame):
        """Creates radio buttons for selection mode."""
        mode_frame = ttk.LabelFrame(parent_frame, text="Selection Mode")
        mode_frame.pack(fill='x', padx=5, pady=5)
        
        ttk.Radiobutton(mode_frame, text="Single File", 
                       variable=self.selection_mode, 
                       value="file").pack(side=tk.LEFT, padx=10)
        ttk.Radiobutton(mode_frame, text="Folder (Recursive)", 
                       variable=self.selection_mode, 
                       value="folder").pack(side=tk.LEFT, padx=10)

    def select_input_paths(self):
        """Opens a dialog to select input files or directory based on mode."""
        if self.selection_mode.get() == "folder":
            path = filedialog.askdirectory(title="Select Input Directory")
            if path:
                self.input_paths = [Path(path)]
                self.update_input_display()
                return True
        else:
            paths = filedialog.askopenfilenames(
                title="Select Input Files",
                filetypes=[("Supported Files", 
                          [f"*{ext}" for ext in self.supported_extensions])]
            )
            if paths:
                self.input_paths = [Path(p) for p in paths]
                self.update_input_display()
                return True
        return False

    def select_output_path(self):
        """Opens a dialog to select an output directory."""
        path = filedialog.askdirectory(title="Select Output Directory")
        if path:
            self.output_path = Path(path)
            logging.info(f"Output path selected: {self.output_path}")
            self.update_output_display()
            return True
        return False

    def set_same_as_input(self):
        """Sets the output path to be the same as the input path."""
        if self.input_paths:
            if self.selection_mode.get() == "folder":
                self.output_path = self.input_paths[0]
            else:
                self.output_path = self.input_paths[0].parent
            logging.info(f"Output path set to same as input: {self.output_path}")
            self.update_output_display()
            return True
        else:
            messagebox.showwarning("Warning", "Please select input paths first.")
            return False

    def update_input_display(self):
        """Updates the input paths display."""
        if hasattr(self, 'input_paths_display'):
            self.input_paths_display.configure(state='normal')
            self.input_paths_display.delete(1.0, tk.END)
            for path in self.input_paths:
                self.input_paths_display.insert(tk.END, f"{path}\n")
            self.input_paths_display.configure(state='disabled')

    def update_output_display(self):
        """Updates the output path display."""
        if hasattr(self, 'output_path_display'):
            self.output_path_display.configure(state='normal')
            self.output_path_display.delete(1.0, tk.END)
            if self.output_path:
                self.output_path_display.insert(tk.END, str(self.output_path))
            self.output_path_display.configure(state='disabled')

    def get_all_files_recursive(self, directory: Path) -> list:
        """Recursively gets all supported files from directory."""
        files = []
        try:
            for item in directory.rglob("*"):
                if item.is_file() and item.suffix.lower() in self.supported_extensions:
                    files.append(item)
        except Exception as e:
            self.send_progress_update(f"Error scanning directory {directory}: {e}")
        return sorted(files)  # Sort files for consistent processing order

    def process_paths(self):
        """Enhanced process_paths method with stop functionality."""
        if not self.input_paths:
            messagebox.showerror("Error", "No input paths selected.")
            return

        if self.output_path is None:
            result = messagebox.askyesno(
                "Question", 
                "No output path selected. Output to same directory as input?"
            )
            if result:
                if not self.set_same_as_input():
                    return
            else:
                messagebox.showinfo("Info", "Processing cancelled.")
                return

        # Reset stop flag before starting new processing
        self.stop_flag.clear()
        
        # Store thread reference for stopping
        self.processing_thread = threading.Thread(
            target=self._process_paths_threaded, 
            daemon=True
        )
        self.processing_thread.start()

    def _process_paths_threaded(self):
        """Enhanced threaded processing with stop functionality."""
        try:
            self.before_processing()
            
            files_to_process = []
            if self.selection_mode.get() == "folder":
                for input_path in self.input_paths:
                    files_to_process.extend(self.get_all_files_recursive(input_path))
            else:
                files_to_process = self.input_paths

            total_files = len(files_to_process)
            self.send_progress_update(f"Found {total_files} files to process")

            for index, file_path in enumerate(files_to_process, 1):
                # Check if processing should stop
                if self.stop_flag.is_set():
                    self.send_progress_update("Processing stopped by user")
                    return

                try:
                    if self.selection_mode.get() == "folder":
                        relative_path = file_path.parent.relative_to(self.input_paths[0])
                        output_dir = self.output_path / relative_path
                    else:
                        output_dir = self.output_path

                    output_dir.mkdir(parents=True, exist_ok=True)

                    self.send_progress_update(
                        f"Processing file {index}/{total_files}: {file_path.name}"
                    )
                    self.process_file(file_path, output_dir)

                except Exception as e:
                    self.send_progress_update(
                        f"Error processing {file_path.name}: {str(e)}"
                    )
                    logging.exception(f"Error processing {file_path}")

            self.after_processing()

        except Exception as e:
            error_msg = f"Error during processing: {str(e)}"
            self.send_progress_update(error_msg)
            logging.exception(error_msg)
        finally:
            self.stop_flag.clear()
            self.processing_thread = None
            self.send_progress_update("Processing complete")


    def send_progress_update(self, message: str):
        """Sends a progress update message to the GUI."""
        self.progress_queue.put(message)

    def before_processing(self):
        """Hook for pre-processing setup."""
        pass

    def after_processing(self):
        """Hook for post-processing cleanup."""
        pass

    def process_file(self, input_file: Path, output_dir: Path):
        """Abstract method for processing a single file."""
        raise NotImplementedError("Subclasses must implement process_file()")
    
    def should_skip_file(self, input_file: Path, output_dir: Path, output_extension: str = None) -> bool:
        """Check if processing should be skipped based on existing output."""
        if not self.check_output_exists.get():
            return False
        
        # If no specific extension provided, cannot check
        if not output_extension:
            return False
            
        # Construct expected output filename
        output_filename = input_file.stem + output_extension
        output_file = output_dir / output_filename
        
        # Use consolidated file_utils function
        return should_skip_processing(input_file, output_file, check_exists=True)

    def handle_drop(self, event):
        """Handles drag and drop events."""
        files = event.data.split()
        
        # Filter files based on selection mode and supported extensions
        if self.selection_mode.get() == "folder":
            self.input_paths = [Path(f) for f in files if Path(f).is_dir()]
        else:
            self.input_paths = [
                Path(f) for f in files 
                if Path(f).is_file() and Path(f).suffix.lower() in self.supported_extensions
            ]
        
        if self.input_paths:
            logging.info(f"Files dropped: {self.input_paths}")
            self.update_input_display()
        else:
            messagebox.showwarning(
                "Warning", 
                "No valid input paths found. Please check selection mode and file types."
            )

class LanguageSelectionMixin:
    """Mixin class for language selection functionality."""
    
    def create_language_selection(self):
        """Creates the language selection UI elements."""
        # Language selection frame
        self.lang_frame = ttk.LabelFrame(self.master, text="Language Selection")
        self.lang_frame.pack(fill='x', padx=5, pady=5)

        # Source language
        source_frame = ttk.Frame(self.lang_frame)
        source_frame.pack(side=tk.LEFT, padx=5, pady=5)
        
        ttk.Label(source_frame, text="Source Language:").pack(side=tk.LEFT, padx=5)
        self.source_lang_combo = ttk.Combobox(
            source_frame, 
            textvariable=self.source_lang,
            values=list(self.supported_languages.get("source_languages", {}).keys()),
            state="readonly",
            width=10
        )
        self.source_lang_combo.pack(side=tk.LEFT, padx=5)

        # Target language
        target_frame = ttk.Frame(self.lang_frame)
        target_frame.pack(side=tk.LEFT, padx=5, pady=5)
        
        ttk.Label(target_frame, text="Target Language:").pack(side=tk.LEFT, padx=5)
        self.target_lang_combo = ttk.Combobox(
            target_frame, 
            textvariable=self.target_lang,
            values=list(self.supported_languages.get("target_languages", {}).keys()),
            state="readonly",
            width=10
        )
        self.target_lang_combo.pack(side=tk.LEFT, padx=5)

        # Add language info tooltips
        self.add_language_tooltips()

    def add_language_tooltips(self):
        """Adds tooltips to show full language names on hover."""
        def create_tooltip(widget, text):
            def show_tooltip(event):
                tooltip = tk.Toplevel()
                tooltip.wm_overrideredirect(True)
                tooltip.wm_geometry(f"+{event.x_root+10}+{event.y_root+10}")
                
                label = ttk.Label(tooltip, text=text, background="#ffffe0", 
                                relief='solid', borderwidth=1)
                label.pack()
                
                def hide_tooltip():
                    tooltip.destroy()
                
                widget.tooltip = tooltip
                widget.after(2000, hide_tooltip)
            
            def hide_tooltip(event):
                if hasattr(widget, 'tooltip'):
                    widget.tooltip.destroy()
            
            widget.bind('<Enter>', show_tooltip)
            widget.bind('<Leave>', hide_tooltip)

        # Add tooltips for source languages
        source_languages = self.supported_languages.get("source_languages", {})
        source_tooltip_text = "Available source languages:\n" + \
                            "\n".join(f"{code}: {name}" for code, name in source_languages.items())
        create_tooltip(self.source_lang_combo, source_tooltip_text)

        # Add tooltips for target languages
        target_languages = self.supported_languages.get("target_languages", {})
        target_tooltip_text = "Available target languages:\n" + \
                            "\n".join(f"{code}: {name}" for code, name in target_languages.items())
        create_tooltip(self.target_lang_combo, target_tooltip_text)


class TextToSpeechTool(ToolBase):
    """
    Converts text files (.txt) to MP3 audio using the ElevenLabs API.
    Uses TextToSpeechCore for the actual processing with filename-based voice selection.
    Voice name can be separated by underscore, hyphen, or space.
    """

    def __init__(self, master, config_manager, progress_queue):
        super().__init__(master, config_manager, progress_queue)
        self.supported_extensions = {'.txt'}
        self.config_manager = config_manager

    def before_processing(self):
        """Pre-processing setup."""
        api_key = self.config_manager.get_api_keys().get("elevenlabs")
        if not api_key:
            raise ValueError("ElevenLabs API key not configured. Please add your API key in the Configuration menu.")

    def process_file(self, input_file: Path, output_dir: Path):
        """Processes a single text file for TTS conversion using consolidated audio processor."""
        try:
            # Check for interruption
            if self.stop_flag.is_set():
                raise InterruptedError("Processing stopped by user")
            
            # Determine output path
            output_file = output_dir / f"{input_file.stem}.mp3"
            
            # Create audio processor with GUI-specific progress reporter
            progress_reporter = ProgressReporter(callback=lambda msg: (
                self.send_progress_update(msg) if not self.stop_flag.is_set() 
                else (_ for _ in ()).throw(InterruptedError("Processing stopped by user"))
            ))
            
            processor = create_audio_processor(
                self.service_manager,
                progress_reporter=progress_reporter,
                config=ProcessorConfig(
                    skip_existing=self.check_output_exists.get(),
                    allowed_extensions={'.txt'},
                    max_file_size_mb=10.0
                )
            )
            
            # Process the file using consolidated processor
            result = processor.process_file(
                input_file,
                output_file,
                operation='synthesize'
            )
            
            if result.success:
                self.send_progress_update(f"Successfully generated audio: {output_file.name}")
            elif result.skipped:
                # Skip message already sent by processor
                pass
            else:
                self.send_progress_update(f"Failed to generate audio for: {input_file.name} - {result.message}")
            
        except InterruptedError:
            raise
        except Exception as e:
            error_message = f"Error generating audio for {input_file.name}: {str(e)}"
            self.send_progress_update(error_message)
            logging.exception(error_message)

class PPTXTranslationTool(ToolBase, LanguageSelectionMixin):
    """Translates PPTX files from one language to another."""

    def __init__(self, master, config_manager, progress_queue):
        super().__init__(master, config_manager, progress_queue)
        self.supported_extensions = {'.pptx'}
        
        # Language selection variables
        self.source_lang = tk.StringVar(value="en")
        self.target_lang = tk.StringVar(value="fr")
        
        self.api_key = self.config_manager.get_api_keys().get("deepl")
        if not self.api_key:
            logging.warning("DeepL API key not configured")

    def create_language_selection(self):
        """Creates the language selection UI elements."""
        # Language selection frame
        self.lang_frame = ttk.LabelFrame(self.master, text="Language Selection")
        self.lang_frame.pack(fill='x', padx=5, pady=5)

        # Source language
        source_frame = ttk.Frame(self.lang_frame)
        source_frame.pack(side=tk.LEFT, padx=5, pady=5)
        
        ttk.Label(source_frame, text="Source Language:").pack(side=tk.LEFT, padx=5)
        self.source_lang_combo = ttk.Combobox(
            source_frame, 
            textvariable=self.source_lang,
            values=list(self.supported_languages.get("source_languages", {}).keys()),
            state="readonly",
            width=10
        )
        self.source_lang_combo.pack(side=tk.LEFT, padx=5)

        # Target language
        target_frame = ttk.Frame(self.lang_frame)
        target_frame.pack(side=tk.LEFT, padx=5, pady=5)
        
        ttk.Label(target_frame, text="Target Language:").pack(side=tk.LEFT, padx=5)
        self.target_lang_combo = ttk.Combobox(
            target_frame, 
            textvariable=self.target_lang,
            values=list(self.supported_languages.get("target_languages", {}).keys()),
            state="readonly",
            width=10
        )
        self.target_lang_combo.pack(side=tk.LEFT, padx=5)

        # Add language info tooltips
        self.add_language_tooltips()

    def add_language_tooltips(self):
        """Adds tooltips to show full language names on hover."""
        def create_tooltip(widget, text):
            def show_tooltip(event):
                tooltip = tk.Toplevel()
                tooltip.wm_overrideredirect(True)
                tooltip.wm_geometry(f"+{event.x_root+10}+{event.y_root+10}")
                
                label = ttk.Label(tooltip, text=text, background="#ffffe0", relief='solid', borderwidth=1)
                label.pack()
                
                def hide_tooltip():
                    tooltip.destroy()
                
                widget.tooltip = tooltip
                widget.after(2000, hide_tooltip)
            
            def hide_tooltip(event):
                if hasattr(widget, 'tooltip'):
                    widget.tooltip.destroy()
            
            widget.bind('<Enter>', show_tooltip)
            widget.bind('<Leave>', hide_tooltip)

        # Add tooltips for source languages
        source_languages = self.supported_languages.get("source_languages", {})
        source_tooltip_text = "Available source languages:\n" + \
                            "\n".join(f"{code}: {name}" for code, name in source_languages.items())
        create_tooltip(self.source_lang_combo, source_tooltip_text)

        # Add tooltips for target languages
        target_languages = self.supported_languages.get("target_languages", {})
        target_tooltip_text = "Available target languages:\n" + \
                            "\n".join(f"{code}: {name}" for code, name in target_languages.items())
        create_tooltip(self.target_lang_combo, target_tooltip_text)



    def process_file(self, input_file: Path, output_dir: Path = None):
        """Processes a single PPTX file using consolidated translation processor."""
        if input_file.suffix.lower() != ".pptx":
            self.send_progress_update(f"Skipping non-PPTX file: {input_file}")
            return

        try:
            # Check for interruption
            if self.stop_flag.is_set():
                raise InterruptedError("Processing stopped by user")
                
            # Create output filename with target language suffix
            output_filename = f"{input_file.stem}_{self.target_lang.get()}{input_file.suffix}"
            output_file = output_dir / output_filename
            
            # Create translation processor with GUI-specific progress reporter
            progress_reporter = ProgressReporter(callback=lambda msg: (
                self.send_progress_update(msg) if not self.stop_flag.is_set() 
                else (_ for _ in ()).throw(InterruptedError("Processing stopped by user"))
            ))
            
            processor = create_translation_processor(
                self.service_manager,
                progress_reporter=progress_reporter,
                config=ProcessorConfig(
                    skip_existing=self.check_output_exists.get(),
                    allowed_extensions={'.pptx'},
                    max_file_size_mb=50.0
                )
            )
            
            # Process the file using consolidated processor
            result = processor.process_file(
                input_file,
                output_file,
                source_language=self.source_lang.get(),
                target_language=self.target_lang.get()
            )
            
            if result.success:
                self.send_progress_update(f"Successfully translated: {output_file.name}")
            elif result.skipped:
                # Skip message already sent by processor
                pass
            else:
                self.send_progress_update(f"Failed to translate: {input_file.name} - {result.message}")

        except InterruptedError:
            raise
        except Exception as e:
            error_message = f"Error translating {input_file.name}: {e}"
            self.send_progress_update(error_message)
            logging.exception(error_message)




class AudioTranscriptionTool(ToolBase):
    """Transcribes audio files using OpenAI's Whisper API via AudioTranscriptionCore."""

    def __init__(self, master, config_manager, progress_queue):
        super().__init__(master, config_manager, progress_queue)
        # Define supported extensions
        self.supported_extensions = {'.wav', '.mp3', '.m4a', '.webm', 
                                   '.mp4', '.mpga', '.mpeg'}

        # Get OpenAI API key from config
        self.api_key = self.config_manager.get_api_keys().get("openai")
        if not self.api_key:
            logging.warning("OpenAI API key not configured")

    def process_file(self, input_file: Path, output_dir: Path):
        """Processes a single audio file using consolidated audio processor."""
        try:
            # Check if processing should stop
            if self.stop_flag.is_set():
                raise InterruptedError("Processing stopped by user")
            
            # Create output filename with transcript suffix
            output_filename = f"{input_file.stem}_transcript.txt"
            output_path = output_dir / output_filename
            
            # Create audio processor with GUI-specific progress reporter
            progress_reporter = ProgressReporter(callback=lambda msg: (
                self.send_progress_update(msg) if not self.stop_flag.is_set() 
                else (_ for _ in ()).throw(InterruptedError("Processing stopped by user"))
            ))
            
            processor = create_audio_processor(
                self.service_manager,
                progress_reporter=progress_reporter,
                config=ProcessorConfig(
                    skip_existing=self.check_output_exists.get(),
                    allowed_extensions=self.supported_extensions,
                    max_file_size_mb=100.0
                )
            )
            
            # Process the file using consolidated processor
            result = processor.process_file(
                input_file,
                output_path,
                operation='transcribe'
            )
            
            if result.success:
                self.send_progress_update(f"Successfully transcribed: {input_file.name}")
            elif result.skipped:
                # Skip message already sent by processor
                pass
            else:
                self.send_progress_update(f"Failed to transcribe: {input_file.name} - {result.message}")

        except InterruptedError:
            raise
        except Exception as e:
            error_message = f"Error transcribing {input_file.name}: {e}"
            self.send_progress_update(error_message)
            logging.exception(error_message)

class TextTranslationTool(ToolBase, LanguageSelectionMixin):
    """Translates text files using DeepL API."""

    def __init__(self, master, config_manager, progress_queue):
        super().__init__(master, config_manager, progress_queue)
        
        # Define supported extensions
        self.supported_extensions = {'.txt'}
        
        # Language selection variables
        self.source_lang = tk.StringVar(value="en")
        self.target_lang = tk.StringVar(value="fr")
        
        # Get API key
        self.api_key = self.config_manager.get_api_keys().get("deepl")
        if not self.api_key:
            logging.warning("DeepL API key not configured")

    def create_language_selection(self):
        """Creates the language selection UI elements."""
        # Language selection frame
        self.lang_frame = ttk.LabelFrame(self.master, text="Language Selection")
        self.lang_frame.pack(fill='x', padx=5, pady=5)

        # Source language
        source_frame = ttk.Frame(self.lang_frame)
        source_frame.pack(side=tk.LEFT, padx=5, pady=5)
        
        ttk.Label(source_frame, text="Source Language:").pack(side=tk.LEFT, padx=5)
        self.source_lang_combo = ttk.Combobox(
            source_frame, 
            textvariable=self.source_lang,
            values=list(self.supported_languages.get("source_languages", {}).keys()),
            state="readonly",
            width=10
        )
        self.source_lang_combo.pack(side=tk.LEFT, padx=5)

        # Target language
        target_frame = ttk.Frame(self.lang_frame)
        target_frame.pack(side=tk.LEFT, padx=5, pady=5)
        
        ttk.Label(target_frame, text="Target Language:").pack(side=tk.LEFT, padx=5)
        self.target_lang_combo = ttk.Combobox(
            target_frame, 
            textvariable=self.target_lang,
            values=list(self.supported_languages.get("target_languages", {}).keys()),
            state="readonly",
            width=10
        )
        self.target_lang_combo.pack(side=tk.LEFT, padx=5)

    def process_file(self, input_file: Path, output_dir: Path):
        """Processes a single text file using consolidated translation processor."""
        try:
            # Check for interruption
            if self.stop_flag.is_set():
                raise InterruptedError("Processing stopped by user")

            # Create output filename with target language suffix
            output_filename = f"{input_file.stem}_{self.target_lang.get()}{input_file.suffix}"
            output_file = output_dir / output_filename
            
            # Create translation processor with GUI-specific progress reporter
            progress_reporter = ProgressReporter(callback=lambda msg: (
                self.send_progress_update(msg) if not self.stop_flag.is_set() 
                else (_ for _ in ()).throw(InterruptedError("Processing stopped by user"))
            ))
            
            processor = create_translation_processor(
                self.service_manager, 
                progress_reporter=progress_reporter,
                config=ProcessorConfig(
                    skip_existing=self.check_output_exists.get(),
                    allowed_extensions={'.txt'},
                    max_file_size_mb=50.0
                )
            )
            
            # Process the file using consolidated processor
            result = processor.process_file(
                input_file, 
                output_file,
                source_language=self.source_lang.get(),
                target_language=self.target_lang.get()
            )
            
            if result.success:
                self.send_progress_update(f"Successfully translated: {input_file.name}")
            elif result.skipped:
                # Skip message already sent by processor
                pass
            else:
                self.send_progress_update(f"Failed to translate: {input_file.name} - {result.message}")

        except InterruptedError:
            raise
        except Exception as e:
            error_message = f"Error translating {input_file.name}: {str(e)}"
            self.send_progress_update(error_message)
            logging.exception(error_message)

    def before_processing(self):
        """Pre-processing setup."""
        if not self.api_key:
            raise ValueError("DeepL API key not configured. Please add your API key in the Configuration menu.")

    def after_processing(self):
        """Post-processing cleanup."""
        pass


class PPTXtoPDFTool(ToolBase):
    """
    Converts PPTX files to PDF, PNG, or WEBP using ConvertAPI service.
    Supports PDF conversion, PNG conversion (one PNG per slide), and WEBP conversion (PNG to WEBP).
    """

    def __init__(self, master, config_manager, progress_queue):
        super().__init__(master, config_manager, progress_queue)
        self.supported_extensions = {'.pptx', '.ppt', '.potx', '.pps', '.ppsx'}
        self.output_format = tk.StringVar(value="pdf")  # Default to PDF

    def create_specific_controls(self, parent_frame):
        """Creates UI elements specific to this tool (output format selection)."""
        format_frame = ttk.LabelFrame(parent_frame, text="Output Format")
        format_frame.pack(fill='x', padx=5, pady=5)

        ttk.Radiobutton(format_frame, text="PDF",
                       variable=self.output_format,
                       value="pdf").pack(side=tk.LEFT, padx=10)
        ttk.Radiobutton(format_frame, text="PNG (One per slide)",
                       variable=self.output_format,
                       value="png").pack(side=tk.LEFT, padx=10)
        ttk.Radiobutton(format_frame, text="WEBP (One per slide)",
                       variable=self.output_format,
                       value="webp").pack(side=tk.LEFT, padx=10)

    def before_processing(self):
        """Load credentials before starting the batch."""
        api_keys = self.config_manager.get_api_keys()
        api_secret = api_keys.get("convertapi")
        
        if not api_secret:
            raise ValueError("Missing ConvertAPI secret in configuration. Please configure it via the menu.")

    def process_file(self, input_file: Path, output_dir: Path):
        """
        Processes a single PPTX file using consolidated conversion processor.
        Converts to PDF, PNG, or WEBP with proper sequential naming.
        """
        try:
            # Check for interruption
            if self.stop_flag.is_set():
                raise InterruptedError("Processing stopped by user")
            
            output_format = self.output_format.get()
            
            # Determine output path based on format
            if output_format == 'pdf':
                output_file = output_dir / f"{input_file.stem}.pdf"
            else:
                # For image formats, the processor will handle multiple output files
                output_file = output_dir / f"{input_file.stem}_slide_01.{output_format}"
            
            # Create conversion processor with GUI-specific progress reporter
            progress_reporter = ProgressReporter(callback=lambda msg: (
                self.send_progress_update(msg) if not self.stop_flag.is_set() 
                else (_ for _ in ()).throw(InterruptedError("Processing stopped by user"))
            ))
            
            processor = create_conversion_processor(
                self.service_manager,
                progress_reporter=progress_reporter,
                config=ProcessorConfig(
                    skip_existing=self.check_output_exists.get(),
                    allowed_extensions=self.supported_extensions,
                    output_formats=['pdf', 'png', 'webp'],
                    max_file_size_mb=25.0
                )
            )
            
            # Process the file using consolidated processor
            result = processor.process_file(
                input_file,
                output_file,
                output_format=output_format
            )
            
            if result.success:
                self.send_progress_update(f"Successfully converted: {input_file.name}")
                return True
            elif result.skipped:
                # Skip message already sent by processor
                return True
            else:
                self.send_progress_update(f"Failed to convert: {input_file.name} - {result.message}")
                return False

        except InterruptedError:
            raise
        except Exception as e:
            error_msg = f"Failed to convert {input_file.name}: {str(e)}"
            self.send_progress_update(f"ERROR: {error_msg}")
            logging.error(error_msg, exc_info=True)
            return False

    def after_processing(self):
        """Cleanup after processing batch."""
        self.send_progress_update("PPTX conversion batch finished.")

class VideoMergeTool(ToolBase):
    """
    Merges MP3 audio files with PNG images to create MP4 videos.
    Uses ffmpeg directly instead of moviepy for better reliability.
    
    Matches files based on 2-digit number patterns in filenames,
    where digits are separated by underscore or hyphen.
    Adds a 0.5s silence between clips in the final video.
    """

    def __init__(self, master, config_manager, progress_queue):
        super().__init__(master, config_manager, progress_queue)
        # We don't use self.supported_extensions in the usual way
        self.supported_extensions = {'.mp3', '.png'}
        
        # Force selection mode to folder only
        self.selection_mode = tk.StringVar(value="folder")
        # Toggle for recursive batch mode (process subfolders)
        self.recursive_mode = tk.BooleanVar(value=False)
        
        # Intro/Outro options
        self.use_intro = tk.BooleanVar(value=False)
        self.use_outro = tk.BooleanVar(value=False)
        
        # Paths to intro/outro media files
        self.intro_path = Path(__file__).parent / "media" / "planB_intro.mp4"
        self.outro_path = Path(__file__).parent / "media" / "pbn_outro.mp3"
        
        # Check dependencies
        self._check_dependencies()

    def _check_dependencies(self):
        """Check if ffmpeg is installed and available."""
        self.dependencies_met = False
        try:
            # Check if ffmpeg is available in PATH
            result = subprocess.run(['ffmpeg', '-version'], 
                                   stdout=subprocess.PIPE, 
                                   stderr=subprocess.PIPE, 
                                   text=True)
            if result.returncode == 0:
                self.dependencies_met = True
                version_info = result.stdout.split('\n')[0]
                logging.info(f"Found ffmpeg: {version_info}")
            else:
                logging.error("ffmpeg command returned non-zero exit code")
                self.send_progress_update("ERROR: ffmpeg command failed")
        except FileNotFoundError:
            error_msg = "ffmpeg not found in PATH. Please install ffmpeg first."
            logging.error(error_msg)
            self.send_progress_update(f"ERROR: {error_msg}")
        except Exception as e:
            logging.error(f"Error checking ffmpeg: {str(e)}")
            self.send_progress_update(f"ERROR: Failed to check ffmpeg: {str(e)}")
            
    def create_selection_mode_controls(self, parent_frame):
        """Override to only show folder selection mode."""
        mode_frame = ttk.LabelFrame(parent_frame, text="Input Mode")
        mode_frame.pack(fill='x', padx=5, pady=5)
        
        ttk.Label(mode_frame, text="This tool accepts folders containing matching MP3 and PNG files").pack(padx=10, pady=5)
        note_frame = ttk.Frame(mode_frame)
        note_frame.pack(fill='x', padx=10, pady=5)
        ttk.Label(note_frame, text="Note: Files are matched by the same 2-digit number in their names").pack(side=tk.LEFT)
        # Recursive batch mode toggle: when enabled, each subfolder is processed separately
        ttk.Checkbutton(mode_frame,
                        text="Recursive batch mode (process subfolders)",
                        variable=self.recursive_mode).pack(anchor='w', padx=10, pady=5)
        
    def create_specific_controls(self, parent_frame):
        """Creates UI elements specific to VideoMergeTool (intro/outro options)."""
        options_frame = ttk.LabelFrame(parent_frame, text="Video Options")
        options_frame.pack(fill='x', padx=5, pady=5)
        
        # Intro option
        intro_check = ttk.Checkbutton(options_frame,
                                     text="Add Plan B intro to beginning",
                                     variable=self.use_intro)
        intro_check.pack(anchor='w', padx=10, pady=5)
        
        # Outro option  
        outro_check = ttk.Checkbutton(options_frame,
                                     text="Use PBN outro audio for last slide",
                                     variable=self.use_outro)
        outro_check.pack(anchor='w', padx=10, pady=5)
        
        # Note about outro behavior
        outro_note = ttk.Label(options_frame, 
                              text="Note: When outro is enabled, the last slide will use the outro audio instead of its MP3 file",
                              font=('TkDefaultFont', 9, 'italic'))
        outro_note.pack(anchor='w', padx=25, pady=(0, 5))
        
    def select_input_paths(self):
        """Override to only allow folder selection."""
        path = filedialog.askdirectory(title="Select Folder with MP3 and PNG Files")
        if path:
            self.input_paths = [Path(path)]
            self.update_input_display()
            return True
        return False
            
    def before_processing(self):
        """Pre-processing setup and validation."""
        if not self.dependencies_met:
            self._check_dependencies()  # Check again in case user installed ffmpeg
            if not self.dependencies_met:
                raise ImportError("ffmpeg not found. Please install ffmpeg and make sure it's in your PATH.")

    def process_file(self, input_file, output_dir):
        """Not used for this tool - we process directories instead."""
        pass

    def _process_paths_threaded(self):
        """Override to implement directory-based processing instead of file-based."""
        try:
            self.before_processing()
            
            if not self.input_paths:
                self.send_progress_update("No input directory selected.")
                return
                
            input_dir = self.input_paths[0]
            if not input_dir.is_dir():
                self.send_progress_update(f"Error: {input_dir} is not a directory.")
                return
                
            # Create output directory if it doesn't exist
            output_dir = self.output_path
            output_dir.mkdir(parents=True, exist_ok=True)
                
            # Process the directory
            self.process_directory(input_dir, output_dir)
                
            self.after_processing()
            
        except Exception as e:
            error_msg = f"Error during processing: {str(e)}"
            self.send_progress_update(error_msg)
            logging.exception(error_msg)
        finally:
            self.stop_flag.clear()
            self.processing_thread = None
            self.send_progress_update("Processing complete")
            
    def process_directory(self, input_dir, output_dir):
        """Processes MP3/PNG pairs in either single-folder or recursive mode to create MP4 videos."""
        from core.video_merger import VideoMergerCore
        
        # Determine mode: flat (single folder) or recursive per subfolder
        if not getattr(self, 'recursive_mode', None) or not self.recursive_mode.get():
            self.send_progress_update(f"Processing single-folder mode: {input_dir}")
            # Collect MP3 and PNG files directly in the selected folder
            try:
                entries = os.listdir(input_dir)
            except Exception as e:
                self.send_progress_update(f"Error reading directory {input_dir}: {e}")
                return
            mp3_files = sorted([input_dir / f for f in entries if f.lower().endswith('.mp3')])
            png_files = sorted([input_dir / f for f in entries if f.lower().endswith('.png')])
            self.send_progress_update(f"Found {len(mp3_files)} MP3 and {len(png_files)} PNG in {input_dir}")
            # Match files by identifier
            file_pairs = self.match_file_pairs(mp3_files, png_files)
            if not file_pairs:
                self.send_progress_update(f"No matching MP3/PNG pairs found in {input_dir}")
                return
            self.send_progress_update(f"Found {len(file_pairs)} matching pairs in {input_dir}")
            # Sort file pairs by numeric ID
            file_pairs.sort(key=lambda x: x[0])
            # Generate output filename based on the first MP3 file (remove 2-digit identifier)
            _, first_mp3, _ = file_pairs[0]
            mp3_stem = first_mp3.stem
            identifier_pattern = r'[_-](\d{2})(?:[_-])'
            output_name = re.sub(identifier_pattern, '_', mp3_stem)
            end_pattern = r'[_-]\d{2}$'
            output_name = re.sub(end_pattern, '', output_name)
            # Ensure output directory exists
            output_dir.mkdir(parents=True, exist_ok=True)
            output_file = output_dir / f"{output_name}.mp4"
            self.send_progress_update(f"Output file will be: {output_file}")
            self.create_video_with_ffmpeg(file_pairs, output_file)
            return
        # Recursive mode: walk through all subdirectories
        self.send_progress_update(f"Scanning directory recursively: {input_dir}")
        for dirpath, dirnames, filenames in os.walk(input_dir):
            curr_dir = Path(dirpath)
            # Collect MP3 and PNG files in the current directory
            mp3_files = sorted([curr_dir / f for f in filenames if f.lower().endswith('.mp3')])
            png_files = sorted([curr_dir / f for f in filenames if f.lower().endswith('.png')])
            if not mp3_files or not png_files:
                continue
            self.send_progress_update(f"Found {len(mp3_files)} MP3 and {len(png_files)} PNG in {curr_dir}")
            # Match files by identifier
            file_pairs = self.match_file_pairs(mp3_files, png_files)
            if not file_pairs:
                self.send_progress_update(f"No matching MP3/PNG pairs found in {curr_dir}")
                continue
            self.send_progress_update(f"Found {len(file_pairs)} matching pairs in {curr_dir}")
            # Sort file pairs by numeric ID
            file_pairs.sort(key=lambda x: x[0])
            # Generate output filename based on the first MP3 file (remove 2-digit identifier)
            _, first_mp3, _ = file_pairs[0]
            mp3_stem = first_mp3.stem
            identifier_pattern = r'[_-](\d{2})(?:[_-])'
            output_name = re.sub(identifier_pattern, '_', mp3_stem)
            end_pattern = r'[_-]\d{2}$'
            output_name = re.sub(end_pattern, '', output_name)
            output_file = curr_dir / f"{output_name}.mp4"
            self.send_progress_update(f"Output file will be: {output_file}")
            # Create the output video
            self.create_video_with_ffmpeg(file_pairs, output_file)

        
    def match_file_pairs(self, mp3_files, png_files):
        """
        Match MP3 and PNG files based on a generic two-digit index in their filenames.
        The index is defined as exactly two digits not part of a larger digit sequence
        and can be surrounded by any non-digit character (e.g., '_01_', '-01-', '.01.').
        """
        file_pairs = []
        mp3_dict = {}
        png_dict = {}

        # Compile a regex to match exactly two digits not adjacent to other digits
        id_pattern = re.compile(r'(?<!\d)(\d{2})(?!\d)')

        # Extract indices for PNG files
        for png_file in png_files:
            match = id_pattern.search(png_file.name)
            if match:
                idx = match.group(1)
                self.send_progress_update(f"PNG found index {idx} in {png_file.name}")
                png_dict[idx] = png_file

        # Extract indices for MP3 files
        for mp3_file in mp3_files:
            match = id_pattern.search(mp3_file.name)
            if match:
                idx = match.group(1)
                self.send_progress_update(f"MP3 found index {idx} in {mp3_file.name}")
                mp3_dict[idx] = mp3_file

        # Match pairs by index
        matched_png_indices = set()
        
        for idx in sorted(mp3_dict.keys(), key=lambda x: int(x)):
            mp3_file = mp3_dict[idx]
            png_file = png_dict.get(idx)
            if png_file:
                self.send_progress_update(f"Matched index {idx}: {mp3_file.name} + {png_file.name}")
                file_pairs.append((idx, mp3_file, png_file))
                matched_png_indices.add(idx)
            else:
                self.send_progress_update(f"No PNG match for MP3 index {idx}: {mp3_file.name}")

        # If outro is enabled, also include PNG files without MP3 matches
        # These will use the outro audio
        if self.use_outro.get():
            for idx in sorted(png_dict.keys(), key=lambda x: int(x)):
                if idx not in matched_png_indices:
                    png_file = png_dict[idx]
                    self.send_progress_update(f"PNG without MP3 match (will use outro): {png_file.name}")
                    # Use None as placeholder for MP3 file
                    file_pairs.append((idx, None, png_file))

        # Return pairs sorted by numeric index
        return sorted(file_pairs, key=lambda x: int(x[0]))

    
    def create_video_with_ffmpeg(self, file_pairs, output_file):
        """
        Create a video from matched MP3/PNG pairs using the VideoMergerCore.
        Handles adding silence between clips and intro/outro.
        """
        from core.video_merger import VideoMergerCore
        
        try:
            # Check if should skip
            if self.check_output_exists.get() and output_file.exists():
                self.send_progress_update(f"Skipping video creation - output already exists: {output_file.name}")
                return
                
            # Create progress callback for the merger
            def progress_callback(message):
                if not self.stop_flag.is_set():
                    self.send_progress_update(message)
                else:
                    raise InterruptedError("Processing stopped by user")
            
            # Initialize video merger
            merger = VideoMergerCore(progress_callback)
            
            # Prepare intro/outro paths
            intro_video = self.intro_path if self.use_intro.get() and self.intro_path.exists() else None
            outro_audio = self.outro_path if self.use_outro.get() and self.outro_path.exists() else None
            
            # Use the new create_video_from_file_pairs method
            success = merger.create_video_from_file_pairs(
                file_pairs=file_pairs,
                output_path=output_file,
                silence_duration=0.2,  # 0.2 seconds silence between clips
                intro_video=intro_video,
                outro_audio=outro_audio
            )
            
            if not success:
                raise RuntimeError("Video creation failed")
            
        except Exception as e:
            error_msg = f"Error creating video: {str(e)}"
            self.send_progress_update(error_msg)
            logging.exception(error_msg)



class SequentialProcessingTool(ToolBase, LanguageSelectionMixin):
    def __init__(self, master, config_manager, progress_queue):
        super().__init__(master, config_manager, progress_queue)
        self.supported_extensions = {'.pptx'}
        
        # Language selection variables
        self.source_lang = tk.StringVar(value="en")
        self.target_lang = tk.StringVar(value="fr")
        
        # Multiple target languages selection
        self.selected_target_langs = set()
        
        # Force selection mode to folder
        self.selection_mode = tk.StringVar(value="folder")
        
        # Initialize sub-tools
        self.pptx_translation_tool = PPTXTranslationTool(master, config_manager, progress_queue)
        self.pptx_export_tool = PPTXtoPDFTool(master, config_manager, progress_queue)
        self.text_translation_tool = TextTranslationTool(master, config_manager, progress_queue)
        self.text_to_speech_tool = TextToSpeechTool(master, config_manager, progress_queue)
        self.video_merge_tool = VideoMergeTool(master, config_manager, progress_queue)

    def create_language_selection(self):
        """Creates enhanced language selection UI with multiple target language support."""
        # Main language frame
        self.lang_frame = ttk.LabelFrame(self.master, text="Language Selection")
        self.lang_frame.pack(fill='x', padx=5, pady=5)

        # Source language frame
        source_frame = ttk.Frame(self.lang_frame)
        source_frame.pack(side=tk.LEFT, padx=5, pady=5)
        
        ttk.Label(source_frame, text="Source Language:").pack(side=tk.LEFT, padx=5)
        self.source_lang_combo = ttk.Combobox(
            source_frame, 
            textvariable=self.source_lang,
            values=list(self.supported_languages.get("source_languages", {}).keys()),
            state="readonly",
            width=10
        )
        self.source_lang_combo.pack(side=tk.LEFT, padx=5)

        # Target language selection button
        target_button = ttk.Button(
            self.lang_frame,
            text="Select Target Languages",
            command=self.open_target_language_selector
        )
        target_button.pack(side=tk.LEFT, padx=5, pady=5)

        # Selected languages display
        self.selected_langs_display = tk.Text(
            self.lang_frame,
            height=2,
            width=30,
            wrap=tk.WORD,
            state='disabled'
        )
        self.selected_langs_display.pack(side=tk.LEFT, padx=5, pady=5, fill='x', expand=True)

    def open_target_language_selector(self):
        """Opens a dialog for selecting multiple target languages."""
        selector = tk.Toplevel(self.master)
        selector.title("Select Target Languages")
        selector.geometry("300x400")
        
        # Make dialog modal
        selector.transient(self.master)
        selector.grab_set()
        
        # Create scrollable frame
        main_frame = ttk.Frame(selector)
        main_frame.pack(fill='both', expand=True, padx=5, pady=5)
        
        canvas = tk.Canvas(main_frame)
        scrollbar = ttk.Scrollbar(main_frame, orient="vertical", command=canvas.yview)
        scrollable_frame = ttk.Frame(canvas)

        scrollable_frame.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )

        canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)

        # Variables to store checkbutton states
        checkbutton_vars = {}
        
        # Create checkbuttons for each target language
        target_languages = self.supported_languages.get("target_languages", {})
        for code, name in sorted(target_languages.items()):
            var = tk.BooleanVar(value=code in self.selected_target_langs)
            checkbutton_vars[code] = var
            ttk.Checkbutton(
                scrollable_frame,
                text=f"{code} - {name}",
                variable=var
            ).pack(anchor='w', padx=5, pady=2)

        # Pack scrollbar and canvas
        scrollbar.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True)

        def save_selection():
            self.selected_target_langs = {
                code for code, var in checkbutton_vars.items() if var.get()
            }
            self.update_selected_languages_display()
            selector.destroy()

        # Add buttons
        button_frame = ttk.Frame(selector)
        button_frame.pack(fill='x', padx=5, pady=5)
        
        ttk.Button(
            button_frame,
            text="Select All",
            command=lambda: [var.set(True) for var in checkbutton_vars.values()]
        ).pack(side=tk.LEFT, padx=5)
        
        ttk.Button(
            button_frame,
            text="Clear All",
            command=lambda: [var.set(False) for var in checkbutton_vars.values()]
        ).pack(side=tk.LEFT, padx=5)
        
        ttk.Button(
            button_frame,
            text="Save",
            command=save_selection
        ).pack(side=tk.RIGHT, padx=5)

    def update_selected_languages_display(self):
        """Updates the display of selected target languages."""
        self.selected_langs_display.config(state='normal')
        self.selected_langs_display.delete(1.0, tk.END)
        
        if self.selected_target_langs:
            langs_text = ", ".join(sorted(self.selected_target_langs))
            self.selected_langs_display.insert(tk.END, f"Selected: {langs_text}")
        else:
            self.selected_langs_display.insert(tk.END, "No target languages selected")
        
        self.selected_langs_display.config(state='disabled')

    def process_paths(self):
        """Override to process each target language sequentially."""
        if not self.input_paths:
            messagebox.showerror("Error", "No input paths selected.")
            return

        if not self.selected_target_langs:
            messagebox.showerror("Error", "No target languages selected.")
            return

        if self.output_path is None:
            result = messagebox.askyesno(
                "Question", 
                "No output path selected. Output to same directory as input?"
            )
            if result:
                if not self.set_same_as_input():
                    return
            else:
                messagebox.showinfo("Info", "Processing cancelled.")
                return

        # Reset stop flag before starting new processing
        self.stop_flag.clear()
        
        # Store thread reference for stopping
        self.processing_thread = threading.Thread(
            target=self._process_multiple_languages, 
            daemon=True
        )
        self.processing_thread.start()

    def _process_multiple_languages(self):
        """Process all selected target languages sequentially."""
        try:
            self.before_processing()
            
            for target_lang in sorted(self.selected_target_langs):
                if self.stop_flag.is_set():
                    self.send_progress_update("Processing stopped by user")
                    return
                
                self.send_progress_update(f"\nProcessing target language: {target_lang}")
                self.target_lang.set(target_lang)
                
                # Create language-specific output directory
                lang_output_dir = self.output_path / target_lang
                lang_output_dir.mkdir(parents=True, exist_ok=True)
                
                # Process all files for this language
                if self.selection_mode.get() == "folder":
                    for input_path in self.input_paths:
                        files = self.get_all_files_recursive(input_path)
                        for file_path in files:
                            if self.stop_flag.is_set():
                                return
                            self.process_file(file_path, lang_output_dir)
                else:
                    for file_path in self.input_paths:
                        if self.stop_flag.is_set():
                            return
                        self.process_file(file_path, lang_output_dir)
            
            self.after_processing()
            
        except Exception as e:
            self.send_progress_update(f"Error during processing: {str(e)}")
            logging.exception("Error during multiple language processing")


class RewardEvaluatorTool(ToolBase):
    """
    Unified reward evaluator for both PPTX and TXT files.
    Supports different reward modes: Image PPTX, Video PPTX, and TXT.
    """

    def __init__(self, master, config_manager, progress_queue):
        super().__init__(master, config_manager, progress_queue)
        self.supported_extensions = {'.pptx', '.ppt', '.txt'}
        
        # Additional variables specific to this tool
        self.reward_mode = tk.StringVar(value="image")  # image, video, txt
        self.target_language = tk.StringVar(value="en")
        self.results = []
        
        # Initialize the unified evaluator
        from core.unified_reward_evaluator import UnifiedRewardEvaluator
        self.evaluator = UnifiedRewardEvaluator()

    def create_specific_controls(self, parent_frame):
        """Creates UI elements specific to this tool."""
        
        # Language selection
        lang_frame = ttk.LabelFrame(parent_frame, text="Language Selection")
        lang_frame.pack(fill='x', padx=5, pady=5)
        
        # Target language
        target_label = ttk.Label(lang_frame, text="Target Language:")
        target_label.pack(side=tk.LEFT, padx=5)
        
        languages = list(self.evaluator.language_factors.keys())
        target_combo = ttk.Combobox(lang_frame, textvariable=self.target_language, 
                                   values=languages, state="readonly")
        target_combo.pack(side=tk.LEFT, padx=5)
        
        # Reward mode selection
        mode_frame = ttk.LabelFrame(parent_frame, text="Reward Mode")
        mode_frame.pack(fill='x', padx=5, pady=5)
        
        ttk.Radiobutton(mode_frame, text="Image PPTX (factor 1.5)", 
                       variable=self.reward_mode, 
                       value="image").pack(side=tk.LEFT, padx=10)
        ttk.Radiobutton(mode_frame, text="Video PPTX (factor 1.0)", 
                       variable=self.reward_mode, 
                       value="video").pack(side=tk.LEFT, padx=10)
        ttk.Radiobutton(mode_frame, text="TXT Files", 
                       variable=self.reward_mode, 
                       value="txt").pack(side=tk.LEFT, padx=10)
        
        # Results display frame
        results_frame = ttk.LabelFrame(parent_frame, text="Results")
        results_frame.pack(fill='both', expand=True, padx=5, pady=5)
        
        # Results text widget with scrollbar
        text_frame = ttk.Frame(results_frame)
        text_frame.pack(fill='both', expand=True, padx=5, pady=5)
        
        self.results_text = tk.Text(text_frame, height=10, wrap=tk.WORD)
        results_scrollbar = ttk.Scrollbar(text_frame, orient="vertical", 
                                        command=self.results_text.yview)
        self.results_text.configure(yscrollcommand=results_scrollbar.set)
        
        self.results_text.pack(side=tk.LEFT, fill='both', expand=True)
        results_scrollbar.pack(side=tk.RIGHT, fill='y')
        
        # CSV export button
        export_frame = ttk.Frame(parent_frame)
        export_frame.pack(fill='x', padx=5, pady=5)
        
        self.export_button = ttk.Button(export_frame, text="Export Results to CSV", 
                                       command=self.export_to_csv, state=tk.DISABLED)
        self.export_button.pack(side=tk.LEFT, padx=5)

    def before_processing(self):
        """Setup before processing starts."""
        self.results = []
        self.results_text.configure(state='normal')
        self.results_text.delete(1.0, tk.END)
        self.results_text.configure(state='disabled')
        self.export_button.configure(state=tk.DISABLED)

    def process_file(self, input_file: Path, output_dir: Path):
        """Process a single file for reward evaluation."""
        try:
            self.send_progress_update(f"Evaluating {input_file.name}...")
            
            # Evaluate the file
            result = self.evaluator.evaluate_file(
                str(input_file), 
                self.target_language.get(),
                self.reward_mode.get()
            )
            
            self.results.append(result)
            
            # Display results
            if 'error' in result:
                self.send_progress_update(f"ERROR: {result['error']}")
                self.update_results_display()
                return False
            else:
                # Handle different result formats
                if self.reward_mode.get() == 'txt':
                    reward_euros = result['reward_euros']
                    word_count = result['word_count']
                    self.send_progress_update(f"✅ {input_file.name}")
                    self.send_progress_update(f"   Reward: €{reward_euros:.4f}")
                    self.send_progress_update(f"   Words: {word_count}")
                else:
                    # PPTX mode
                    total_reward = result['total_reward']
                    total_slides = result['total_slides']
                    total_text_boxes = result['total_text_boxes']
                    total_words = result['total_words']
                    self.send_progress_update(f"✅ {input_file.name}")
                    self.send_progress_update(f"   Reward: €{total_reward:.4f}")
                    self.send_progress_update(f"   Slides: {total_slides}, Text boxes: {total_text_boxes}, Words: {total_words}")
                
                self.update_results_display()
                return True
                
        except Exception as e:
            error_msg = f"Failed to evaluate {input_file.name}: {str(e)}"
            self.send_progress_update(f"ERROR: {error_msg}")
            logging.error(error_msg, exc_info=True)
            return False

    def update_results_display(self):
        """Update the results display widget."""
        self.results_text.configure(state='normal')
        self.results_text.delete(1.0, tk.END)
        
        if not self.results:
            self.results_text.insert(tk.END, "No results yet...")
            self.results_text.configure(state='disabled')
            return
        
        # Get summary stats
        summary = self.evaluator.get_summary_stats(self.results)
        
        if 'error' in summary:
            self.results_text.insert(tk.END, f"Error generating summary: {summary['error']}\n")
            self.results_text.configure(state='disabled')
            return
        
        # Header
        file_type = summary.get('file_type', 'Unknown')
        self.results_text.insert(tk.END, f"{file_type} Reward Evaluation Results\n")
        self.results_text.insert(tk.END, f"="*50 + "\n\n")
        self.results_text.insert(tk.END, f"Total Files: {summary['total_files']}\n")
        
        if file_type == 'PPTX':
            self.results_text.insert(tk.END, f"Total Slides: {summary['total_slides']}\n")
            self.results_text.insert(tk.END, f"Total Text Boxes: {summary['total_text_boxes']}\n")
            self.results_text.insert(tk.END, f"Total Words: {summary['total_words']}\n")
        else:
            self.results_text.insert(tk.END, f"Total Words: {summary['total_words']}\n")
            self.results_text.insert(tk.END, f"Avg Words/File: {summary['average_words_per_file']}\n")
        
        self.results_text.insert(tk.END, f"Total Reward: €{summary['total_reward_euros']:.4f}\n")
        self.results_text.insert(tk.END, f"Avg Reward/File: €{summary['average_reward_per_file']:.4f}\n\n")
        
        # Individual results
        for result in self.results:
            if self.reward_mode.get() == 'txt':
                filename = result.get('file_path', 'Unknown')
                if filename != 'Unknown':
                    filename = Path(filename).name
                
                if 'error' in result:
                    self.results_text.insert(tk.END, f"❌ {filename}: {result['error']}\n")
                else:
                    reward = result.get('reward_euros', 0)
                    words = result.get('word_count', 0)
                    difficulty = result.get('difficulty_factor', 1.0)
                    target_lang = result.get('target_language', 'unknown')
                    
                    self.results_text.insert(tk.END, f"✅ {filename}\n")
                    self.results_text.insert(tk.END, f"   Reward: €{reward:.4f}\n")
                    self.results_text.insert(tk.END, f"   Words: {words}\n")
                    self.results_text.insert(tk.END, f"   Target Language: {target_lang}\n")
                    self.results_text.insert(tk.END, f"   Difficulty Factor: {difficulty}\n\n")
            else:
                # PPTX mode
                filename = result.get('filename', 'Unknown')
                if 'error' in result:
                    self.results_text.insert(tk.END, f"❌ {filename}: {result['error']}\n")
                else:
                    reward = result.get('total_reward', 0)
                    slides = result.get('total_slides', 0)
                    text_boxes = result.get('total_text_boxes', 0)
                    words = result.get('total_words', 0)
                    mode = result.get('mode', 'unknown')
                    
                    self.results_text.insert(tk.END, f"✅ {filename}\n")
                    self.results_text.insert(tk.END, f"   Reward: €{reward:.4f}\n")
                    self.results_text.insert(tk.END, f"   Slides: {slides}, Text boxes: {text_boxes}, Words: {words}\n")
                    self.results_text.insert(tk.END, f"   Mode: {mode}\n\n")
        
        self.results_text.configure(state='disabled')
        
        # Enable export button if we have results
        if self.results:
            self.export_button.configure(state=tk.NORMAL)

    def export_to_csv(self):
        """Export results to CSV file."""
        if not self.results:
            messagebox.showwarning("No Results", "No results to export.")
            return
        
        # Ask user for CSV file location
        csv_file = filedialog.asksaveasfilename(
            title="Save Results as CSV",
            defaultextension=".csv",
            filetypes=[("CSV files", "*.csv"), ("All files", "*.*")]
        )
        
        if csv_file:
            try:
                import csv
                with open(csv_file, 'w', newline='', encoding='utf-8') as f:
                    writer = csv.writer(f)
                    
                    if self.reward_mode.get() == 'txt':
                        # TXT CSV format
                        writer.writerow(['File Path', 'Word Count', 'Target Language', 
                                       'Difficulty Factor', 'Euros per Word', 'Reward (Euros)', 'Reward (Cents)', 'Error'])
                        
                        for result in self.results:
                            if 'error' in result:
                                writer.writerow([result.get('file_path', ''), '', '', '', '', '', '', result['error']])
                            else:
                                writer.writerow([
                                    result.get('file_path', ''),
                                    result.get('word_count', 0),
                                    result.get('target_language', ''),
                                    result.get('difficulty_factor', 0),
                                    result.get('euros_per_word', 0),
                                    result.get('reward_euros', 0),
                                    result.get('reward_cents', 0),
                                    ''
                                ])
                    else:
                        # PPTX CSV format
                        writer.writerow(['Filename', 'Total Slides', 'Total Text Boxes', 'Total Words', 
                                       'Mode', 'Total Reward (Euros)', 'Language', 'Error'])
                        
                        for result in self.results:
                            if 'error' in result:
                                writer.writerow([result.get('filename', ''), '', '', '', '', '', '', result['error']])
                            else:
                                writer.writerow([
                                    result.get('filename', ''),
                                    result.get('total_slides', 0),
                                    result.get('total_text_boxes', 0),
                                    result.get('total_words', 0),
                                    result.get('mode', ''),
                                    result.get('total_reward', 0),
                                    result.get('language', ''),
                                    ''
                                ])
                
                self.send_progress_update(f"Results exported to: {csv_file}")
                messagebox.showinfo("Export Successful", f"Results exported to:\n{csv_file}")
            except Exception as e:
                error_msg = f"Failed to export CSV: {str(e)}"
                self.send_progress_update(f"ERROR: {error_msg}")
                messagebox.showerror("Export Error", error_msg)

    def after_processing(self):
        """Cleanup after processing is complete."""
        self.send_progress_update("Reward evaluation completed.")
        self.update_results_display()

    def get_all_files_recursive(self, directory: Path):
        """Get all supported files from directory recursively."""
        files = []
        supported_exts = self.evaluator.get_supported_extensions(self.reward_mode.get())
        for ext in supported_exts:
            files.extend(directory.rglob(f"*{ext}"))
        return files

class TranscriptCleanerTool(ToolBase):
    """Clean and tighten raw transcripts using Claude AI"""
    
    def __init__(self, parent, config_manager, progress_queue):
        super().__init__(parent, config_manager, progress_queue)
        from core.transcript_cleaner import TranscriptCleanerCore
        
        self.supported_extensions = ['.txt']
        self.title = "Clean Raw Transcript"
        self.description = "Clean and tighten raw audio transcripts"
        
        # Get Anthropic API key from config
        api_keys = config_manager.get_api_keys()
        api_key = api_keys.get('anthropic', '')
        
        if api_key:
            self.api_key = api_key
            self.tool_core = TranscriptCleanerCore(
                api_key=api_key,
                progress_callback=self.update_progress
            )
        else:
            self.api_key = None
            self.tool_core = None
    
    def process_file(self, input_path, output_path):
        """Process a single transcript file"""
        if not self.tool_core:
            raise ValueError("Anthropic API key not configured. Please configure API keys first.")
        
        try:
            input_p = Path(input_path)
            
            # For transcript cleaning, output should have -ai-cleaned.txt suffix
            if not output_path:
                output_p = input_p.parent / f"{input_p.stem}-ai-cleaned.txt"
            else:
                output_p = Path(output_path)
                # Ensure output has the correct suffix
                if not output_p.name.endswith('-ai-cleaned.txt'):
                    output_p = output_p.parent / f"{output_p.stem}-ai-cleaned.txt"
            
            # Clean the transcript
            success = self.tool_core.clean_transcript_file(input_p, output_p)
            
            if success:
                self.update_progress(f"✓ Cleaned transcript saved: {output_p.name}")
                return str(output_p)
            else:
                raise Exception("Failed to clean transcript")
                
        except Exception as e:
            error_msg = f"Error cleaning transcript: {str(e)}"
            self.update_progress(error_msg)
            raise Exception(error_msg)
    
    def process_folder(self, folder_path, recursive=False):
        """Process all transcript files in a folder"""
        if not self.tool_core:
            raise ValueError("Anthropic API key not configured. Please configure API keys first.")
        
        try:
            folder_p = Path(folder_path)
            processed_files = self.tool_core.clean_folder(folder_p, recursive=recursive)
            
            if processed_files:
                self.update_progress(f"✓ Successfully cleaned {len(processed_files)} transcripts")
                return processed_files
            else:
                self.update_progress("No transcripts were cleaned")
                return []
                
        except Exception as e:
            error_msg = f"Error processing folder: {str(e)}"
            self.update_progress(error_msg)
            raise Exception(error_msg)


class MainApp(TkinterDnD.Tk):
    """Main application class."""

    def __init__(self):
        super().__init__()

        self.title("Course Video Tools")
        self.geometry("800x600")

        # Initialize components
        self.config_manager = ConfigManager(
            use_project_api_keys=True,
            languages_file=SUPPORTED_LANGUAGES_FILE,
            api_keys_file=API_KEYS_FILE
        )
        self.progress_queue = queue.Queue()

        # Create UI
        self.create_widgets()
        self.process_progress_queue()

    def create_widgets(self):
        """Creates the main UI elements."""
        # Menu
        self.menu_bar = tk.Menu(self)
        
        # Configuration menu
        self.config_menu = tk.Menu(self.menu_bar, tearoff=0)
        self.config_menu.add_command(label="API Keys", command=self.open_api_key_config)
        self.menu_bar.add_cascade(label="Configuration", menu=self.config_menu)
        
        # Help menu
        self.help_menu = tk.Menu(self.menu_bar, tearoff=0)
        self.help_menu.add_command(label="Tool Overview", command=self.show_tool_overview)
        self.help_menu.add_command(label="API Requirements", command=self.show_api_requirements)
        self.help_menu.add_separator()
        self.help_menu.add_command(label="About", command=self.show_about)
        self.menu_bar.add_cascade(label="Help", menu=self.help_menu)
        
        self.config(menu=self.menu_bar)

        # Create main paned window for resizable layout
        self.main_paned = ttk.PanedWindow(self, orient=tk.VERTICAL)
        self.main_paned.pack(expand=True, fill="both", padx=10, pady=10)

        # Top pane: Notebook (Tab Control)
        self.top_frame = ttk.Frame(self.main_paned)
        self.main_paned.add(self.top_frame, weight=3)

        self.notebook = ttk.Notebook(self.top_frame)
        self.notebook.pack(expand=True, fill="both")

        # Tool Frames
        self.pptx_translation_tool = self.create_tool_tab("PPTX Translation", PPTXTranslationTool)
        self.audio_transcription_tool = self.create_tool_tab("Audio Transcription", AudioTranscriptionTool)
        self.text_translation_tool = self.create_tool_tab("Text Translation", TextTranslationTool)
        self.transcript_cleaner_tool = self.create_tool_tab("Clean Transcript", TranscriptCleanerTool)
        self.pptx_to_pdf_tool = self.create_tool_tab("PPTX to PDF/PNG/WEBP", PPTXtoPDFTool)
        self.text_to_speech_tool = self.create_tool_tab("Text to Speech", TextToSpeechTool)  
        self.video_merge_tool = self.create_tool_tab("Video Merge", VideoMergeTool)
        self.sequential_tool = self.create_tool_tab("Sequential Processing", SequentialProcessingTool)
        self.reward_evaluator_tool = self.create_tool_tab("Reward Evaluator", RewardEvaluatorTool)

        # Bottom pane: Progress Text Area
        self.bottom_frame = ttk.Frame(self.main_paned)
        self.main_paned.add(self.bottom_frame, weight=1)

        self.progress_label = tk.Label(self.bottom_frame, text="Progress:")
        self.progress_label.pack(pady=(5, 5))

        # Create frame for progress text and scrollbar
        self.progress_frame = ttk.Frame(self.bottom_frame)
        self.progress_frame.pack(expand=True, fill="both", padx=10, pady=(0, 10))

        self.progress_text = tk.Text(self.progress_frame, state="disabled", wrap=tk.WORD)
        self.progress_text.pack(side=tk.LEFT, expand=True, fill="both")

        self.sb = tk.Scrollbar(self.progress_frame, command=self.progress_text.yview)
        self.sb.pack(side=tk.RIGHT, fill=tk.Y)
        self.progress_text['yscrollcommand'] = self.sb.set

    def create_tool_tab(self, tab_name, tool_class):
        """Creates a tab with the specified tool."""
        frame = ttk.Frame(self.notebook)
        self.notebook.add(frame, text=tab_name)
        
        # Map tool classes to description keys
        tool_description_map = {
            PPTXTranslationTool: "pptx_translation",
            AudioTranscriptionTool: "audio_transcription", 
            TextTranslationTool: "text_translation",
            PPTXtoPDFTool: "pptx_to_pdf_png",
            TextToSpeechTool: "text_to_speech",
            VideoMergeTool: "video_merge",
            SequentialProcessingTool: "sequential_processing",
            RewardEvaluatorTool: "reward_evaluator"
        }
        
        # Get tool description key
        tool_key = tool_description_map.get(tool_class)
        
        # Add description at the top of the tab
        if tool_key:
            self.add_tool_description(frame, tool_key)
        
        tool = tool_class(
            master=frame,
            config_manager=self.config_manager,
            progress_queue=self.progress_queue
        )
        
        self.create_tool_ui(frame, tool)
        return tool

    def add_tool_description(self, frame, tool_key):
        """Add description and help information to the tool tab."""
        try:
            # Get tool information
            tool_info = get_tool_info(tool_key)
            if not tool_info:
                return
            
            desc = tool_info["description"]
            req = tool_info["requirements"]
            tips = tool_info["tips"]
            
            # Create description frame
            desc_frame = ttk.LabelFrame(frame, text=" Tool Information", padding=10)
            desc_frame.pack(fill="x", padx=5, pady=5)
            
            # Main description
            desc_label = tk.Label(desc_frame, text=desc["description"], 
                                font=("Arial", 10, "bold"), fg="navy", wraplength=700)
            desc_label.pack(anchor="w")
            
            # Detailed description
            detail_label = tk.Label(desc_frame, text=desc["details"], 
                                  font=("Arial", 9), fg="gray40", wraplength=700)
            detail_label.pack(anchor="w", pady=(2, 8))
            
            # API requirement info
            if req.get("api_required"):
                api_text = f"🔑 Requires: {req['api_required']} API key"
                if req["api_required"] == "Multiple":
                    api_text = "🔑 Requires: Multiple API keys (see Configuration menu)"
                api_label = tk.Label(desc_frame, text=api_text, 
                                   font=("Arial", 9), fg="red")
                api_label.pack(anchor="w")
            else:
                no_api_label = tk.Label(desc_frame, text="✅ No API key required", 
                                      font=("Arial", 9), fg="green")
                no_api_label.pack(anchor="w")
            
            # Quick tips (collapsible)
            if tips:
                tips_frame = ttk.Frame(desc_frame)
                tips_frame.pack(fill="x", pady=(5, 0))
                
                # Tips toggle button
                self.tips_visible = tk.BooleanVar(value=False)
                tips_btn = tk.Button(tips_frame, text="💡 Show Quick Tips", 
                                   command=lambda: self.toggle_tips(tips_frame, tool_key, tips))
                tips_btn.pack(anchor="w")
                
        except Exception as e:
            # Silently handle any errors in description display
            logging.warning(f"Could not add description for {tool_key}: {e}")

    def add_tab_tooltip(self, tab_name, tool_key):
        """Add tooltip to notebook tab."""
        try:
            short_desc = get_short_description(tool_key)
            if short_desc and short_desc != "Tool description not available":
                # Create a simple tooltip by binding to the tab
                self.create_tooltip_for_tab(tab_name, short_desc)
        except Exception as e:
            logging.warning(f"Could not add tooltip for {tool_key}: {e}")

    def create_tooltip_for_tab(self, tab_name, description):
        """Create a tooltip for a specific tab."""
        def show_tooltip(event):
            # Simple tooltip implementation - could be enhanced
            tooltip = tk.Toplevel()
            tooltip.wm_overrideredirect(True)
            tooltip.wm_geometry(f"+{event.x_root + 10}+{event.y_root + 10}")
            
            label = tk.Label(tooltip, text=description, background="lightyellow", 
                           relief="solid", borderwidth=1, font=("Arial", 9),
                           wraplength=300)
            label.pack()
            
            # Auto-hide after 3 seconds
            tooltip.after(3000, tooltip.destroy)
        
        # This is a simplified tooltip - in a full implementation you'd want
        # to bind to the actual tab widget, but tkinter notebook tabs are complex
        pass

    def toggle_tips(self, parent_frame, tool_key, tips):
        """Toggle display of quick tips."""
        # Remove existing tips if any
        for widget in parent_frame.winfo_children():
            if isinstance(widget, tk.Frame) and hasattr(widget, 'tips_frame'):
                widget.destroy()
                return
        
        # Create tips display
        tips_display = tk.Frame(parent_frame)
        tips_display.tips_frame = True  # Mark as tips frame
        tips_display.pack(fill="x", pady=(5, 0))
        
        tk.Label(tips_display, text="Quick Tips:", font=("Arial", 9, "bold")).pack(anchor="w")
        
        for i, tip in enumerate(tips, 1):
            tip_label = tk.Label(tips_display, text=f"  {i}. {tip}", 
                               font=("Arial", 8), fg="gray30", wraplength=650)
            tip_label.pack(anchor="w", padx=(10, 0))

    def create_tool_ui(self, frame, tool):
        """Creates the UI elements for a specific tool."""
        
        tool.create_selection_mode_controls(frame)

        # If the tool has language selection capability, create the UI
        if isinstance(tool, LanguageSelectionMixin):
            tool.create_language_selection()

        if hasattr(tool, 'create_specific_controls') and callable(tool.create_specific_controls):
            tool.create_specific_controls(frame)
        
        # Input Path Selection
        input_frame = ttk.Frame(frame)
        input_frame.pack(pady=5, fill='x', padx=5)


        input_label = ttk.Label(input_frame, text="Input Path(s):")
        input_label.pack(side=tk.LEFT, padx=5)

        input_button = ttk.Button(input_frame, text="Select Input", command=tool.select_input_paths)
        input_button.pack(side=tk.LEFT, padx=5)

        tool.input_paths_display = tk.Text(input_frame, height=3, width=50, wrap=tk.WORD)
        tool.input_paths_display.pack(side=tk.LEFT, padx=5, fill='x', expand=True)


        
        input_scrollbar = ttk.Scrollbar(input_frame, orient="vertical", 
                                      command=tool.input_paths_display.yview)
        input_scrollbar.pack(side=tk.RIGHT, fill='y')
        tool.input_paths_display.configure(yscrollcommand=input_scrollbar.set)
        tool.input_paths_display.configure(state='disabled')
        
        # Output Path Selection
        output_frame = ttk.Frame(frame)
        output_frame.pack(pady=5, fill='x', padx=5)

        output_label = ttk.Label(output_frame, text="Output Path:")
        output_label.pack(side=tk.LEFT, padx=5)

        output_button = ttk.Button(output_frame, text="Select Output", 
                                 command=tool.select_output_path)
        output_button.pack(side=tk.LEFT, padx=5)

        same_as_input_button = ttk.Button(output_frame, text="Same as Input", 
                                        command=tool.set_same_as_input)
        same_as_input_button.pack(side=tk.LEFT, padx=5)

        tool.output_path_display = tk.Text(output_frame, height=1, width=50)
        tool.output_path_display.pack(side=tk.LEFT, padx=5, fill='x', expand=True)
        tool.output_path_display.configure(state='disabled')

        # Output checking option
        check_frame = ttk.Frame(frame)
        check_frame.pack(pady=5, fill='x', padx=5)
        
        check_output_checkbox = ttk.Checkbutton(
            check_frame,
            text="Skip processing if output already exists",
            variable=tool.check_output_exists
        )
        check_output_checkbox.pack(side=tk.LEFT, padx=5)

        # Button Frame
        button_frame = ttk.Frame(frame)
        button_frame.pack(pady=10)

        # Process Button
        tool.process_button = ttk.Button(
            button_frame, 
            text="Process", 
            command=lambda: self.start_processing(tool)
        )
        tool.process_button.pack(side=tk.LEFT, padx=5)

        # Stop Button
        tool.stop_button = ttk.Button(
            button_frame, 
            text="Stop", 
            command=lambda: self.stop_processing(tool),
            state=tk.DISABLED
        )
        tool.stop_button.pack(side=tk.LEFT, padx=5)

        # Drag and Drop
        input_frame.drop_target_register(DND_FILES)
        input_frame.dnd_bind('<<Drop>>', lambda e: self.on_drop(e, tool))

    def start_processing(self, tool):
        """Starts processing and updates button states."""
        tool.process_button.configure(state=tk.DISABLED)
        tool.stop_button.configure(state=tk.NORMAL)
        tool.process_paths()
        
        # Start checking if processing is complete
        self.check_processing_complete(tool)

    def stop_processing(self, tool):
        """Stops processing and updates button states."""
        tool.stop_processing()
        tool.stop_button.configure(state=tk.DISABLED)
        
        # Wait for processing to actually stop
        if tool.processing_thread:
            tool.processing_thread.join(timeout=1.0)
        
        tool.process_button.configure(state=tk.NORMAL)

    def check_processing_complete(self, tool):
        """Checks if processing is complete and updates button states."""
        if not tool.processing_thread or not tool.processing_thread.is_alive():
            tool.process_button.configure(state=tk.NORMAL)
            tool.stop_button.configure(state=tk.DISABLED)
        else:
            # Check again in 100ms
            self.after(100, lambda: self.check_processing_complete(tool))

    def process_progress_queue(self):
        """Processes messages from the progress queue."""
        try:
            while True:
                message = self.progress_queue.get_nowait()
                self.update_progress_text(message)
                
                # Enable process button if processing is complete
                if message == "Processing complete" or message == "Processing stopped by user":
                    for tool in [self.pptx_translation_tool, self.audio_transcription_tool]:
                        tool.process_button.configure(state=tk.NORMAL)
                        tool.stop_button.configure(state=tk.DISABLED)
                
        except queue.Empty:
            pass
        
        self.after(100, self.process_progress_queue)

    def on_drop(self, event, tool):
        """Handles drag and drop events."""
        files = event.data.split()
        tool.input_paths = [Path(f) for f in files]
        logging.info(f"Files dropped: {tool.input_paths}")
        tool.update_input_display()

    def open_api_key_config(self):
        """Opens a dialog to configure API keys, including ConvertAPI key."""
        api_config_window = tk.Toplevel(self)
        api_config_window.title("API Key Configuration")
        api_config_window.geometry("550x220") # Adjusted for fewer fields

        api_keys = self.config_manager.get_api_keys()
        api_entries = {}

        managed_api_names = [
            "openai",
            "deepl",
            "elevenlabs",
            "convertapi"      # Managed here
        ]

        for api_name in managed_api_names:
            frame = ttk.Frame(api_config_window)
            frame.pack(pady=5, padx=10, fill="x")

            if api_name == "convertapi":
                label_text = "ConvertAPI Key:"
            else:
                 label_text = f"{api_name.replace('_', ' ').title()} API Key:"

            # Set the width when CREATING the Label
            label = ttk.Label(frame, text=label_text, width=20) # Set width here
            # Pack the label WITHOUT the width option
            label.pack(side=tk.LEFT, padx=5, anchor='w') # Removed width=20 from here

            entry = ttk.Entry(frame, width=45)
            if "secret" in api_name.lower():
                entry.config(show="*")
            entry.pack(side=tk.LEFT, expand=True, fill="x", padx=5)

            if api_name in api_keys and api_keys[api_name]:
                entry.insert(0, api_keys[api_name])

            api_entries[api_name] = entry

        save_button = ttk.Button(api_config_window, text="Save",
                               command=lambda: self.save_api_keys(api_entries, api_config_window))
        save_button.pack(pady=15)

    def save_api_keys(self, api_entries, window):
        """Saves the API keys entered in the configuration dialog."""
        api_keys = {}
        for name, entry in api_entries.items():
            api_keys[name] = entry.get().strip()
        
        # Update the config_manager with new keys
        self.config_manager.api_keys = api_keys
        self.config_manager.save_api_keys(api_keys)
        
        # Update the API keys in existing tools
        if hasattr(self, 'text_to_speech_tool'):
            self.text_to_speech_tool.api_key = api_keys.get("elevenlabs")
        if hasattr(self, 'audio_transcription_tool'):
            self.audio_transcription_tool.api_key = api_keys.get("openai")
        if hasattr(self, 'text_translation_tool'):
            self.text_translation_tool.api_key = api_keys.get("deepl")
        if hasattr(self, 'pptx_translation_tool'):
            self.pptx_translation_tool.api_key = api_keys.get("deepl")
        
        window.destroy()
        messagebox.showinfo("Success", "API keys saved successfully")

    def process_progress_queue(self):
        """Processes messages from the progress queue."""
        try:
            while True:
                message = self.progress_queue.get_nowait()
                self.update_progress_text(message)
        except queue.Empty:
            pass

        self.after(100, self.process_progress_queue)

    def update_progress_text(self, message):
        """Updates the progress text area with a new message."""
        self.progress_text.config(state="normal")
        self.progress_text.insert(tk.END, message + "\n")
        self.progress_text.see(tk.END)
        self.progress_text.config(state="disabled")

    def show_tool_overview(self):
        """Show overview of all available tools."""
        try:
            from core.tool_descriptions import get_tool_list_for_gui
            
            # Create overview window
            overview_window = tk.Toplevel(self)
            overview_window.title("Tool Overview")
            overview_window.geometry("700x500")
            overview_window.resizable(True, True)
            
            # Create scrollable frame
            canvas = tk.Canvas(overview_window)
            scrollbar = ttk.Scrollbar(overview_window, orient="vertical", command=canvas.yview)
            scrollable_frame = ttk.Frame(canvas)
            
            scrollable_frame.bind(
                "<Configure>",
                lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
            )
            
            canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
            canvas.configure(yscrollcommand=scrollbar.set)
            
            # Title
            title_label = tk.Label(scrollable_frame, text="Language Toolkit - Tool Overview", 
                                 font=("Arial", 16, "bold"), fg="navy")
            title_label.pack(pady=10)
            
            # Get all tools
            tools = get_tool_list_for_gui()
            
            for tool in tools:
                # Tool frame
                tool_frame = ttk.LabelFrame(scrollable_frame, text=tool["title"], padding=10)
                tool_frame.pack(fill="x", padx=10, pady=5)
                
                # Description
                desc_label = tk.Label(tool_frame, text=tool["description"], 
                                    font=("Arial", 10), wraplength=600)
                desc_label.pack(anchor="w")
                
                # API requirement
                if tool["has_api_requirement"]:
                    api_label = tk.Label(tool_frame, text=f"🔑 Requires: {tool['api_required']} API key", 
                                       font=("Arial", 9), fg="red")
                    api_label.pack(anchor="w", pady=(5, 0))
                else:
                    api_label = tk.Label(tool_frame, text="✅ No API key required", 
                                       font=("Arial", 9), fg="green")
                    api_label.pack(anchor="w", pady=(5, 0))
            
            canvas.pack(side="left", fill="both", expand=True)
            scrollbar.pack(side="right", fill="y")
            
        except Exception as e:
            messagebox.showerror("Error", f"Could not show tool overview: {e}")

    def show_api_requirements(self):
        """Show API key requirements for all tools."""
        try:
            from core.tool_descriptions import get_tool_requirements
            
            requirements = get_tool_requirements()
            
            # Create requirements window
            req_window = tk.Toplevel(self)
            req_window.title("API Key Requirements")
            req_window.geometry("600x400")
            
            # Title
            title_label = tk.Label(req_window, text="API Key Requirements", 
                                 font=("Arial", 16, "bold"), fg="navy")
            title_label.pack(pady=10)
            
            # Instructions
            instructions = tk.Label(req_window, 
                                  text="Configure API keys via Configuration → API Keys menu",
                                  font=("Arial", 10), fg="gray40")
            instructions.pack(pady=(0, 10))
            
            # Scrollable frame for requirements
            canvas = tk.Canvas(req_window)
            scrollbar = ttk.Scrollbar(req_window, orient="vertical", command=canvas.yview)
            scrollable_frame = ttk.Frame(canvas)
            
            scrollable_frame.bind(
                "<Configure>",
                lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
            )
            
            canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
            canvas.configure(yscrollcommand=scrollbar.set)
            
            for tool_name, req in requirements.items():
                tool_frame = ttk.LabelFrame(scrollable_frame, text=tool_name.replace("_", " ").title(), padding=10)
                tool_frame.pack(fill="x", padx=10, pady=5)
                
                if req["api_required"]:
                    req_label = tk.Label(tool_frame, text=f"Required: {req['api_required']}", 
                                       font=("Arial", 10, "bold"), fg="red")
                    req_label.pack(anchor="w")
                    
                    desc_label = tk.Label(tool_frame, text=req["api_description"], 
                                        font=("Arial", 9), fg="gray40")
                    desc_label.pack(anchor="w")
                else:
                    no_req_label = tk.Label(tool_frame, text="No API key required", 
                                          font=("Arial", 10), fg="green")
                    no_req_label.pack(anchor="w")
            
            canvas.pack(side="left", fill="both", expand=True, padx=(10, 0))
            scrollbar.pack(side="right", fill="y", padx=(0, 10))
            
        except Exception as e:
            messagebox.showerror("Error", f"Could not show API requirements: {e}")

    def show_about(self):
        """Show about dialog."""
        about_text = """Language Toolkit
        
A comprehensive suite of tools for language processing, document conversion, and content creation.

Features:
• PowerPoint Translation
• Audio Transcription  
• Text Translation
• Document Conversion
• Text-to-Speech Generation
• Video Creation and Merging
• Sequential Processing Workflows

Built with Python and integrates with leading AI services for professional-quality results."""
        
        messagebox.showinfo("About Language Toolkit", about_text)

if __name__ == "__main__":
    app = MainApp()
    app.mainloop()

