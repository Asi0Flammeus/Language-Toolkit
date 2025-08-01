"""
PPTX Reward Evaluator

This module provides functionality to evaluate the proofreading reward for PPTX files
based on the number of text boxes per slide and various factors like language difficulty.

Formula: reward = delta_slide * (words_slide * diff_lang * euros_per_word + text_box_fee * nbr_text_box)
Where:
- delta = 0 if no words, delta = 1 if words
- euros_per_word = 0.0006 (0.06 cents / word)
- text_box_fee is variable based on presentation type
- image_factor: 1.5 for image presentations, 1.0 for video presentations
"""

import os
import csv
from pathlib import Path
from typing import Dict, List, Tuple, Optional
import yaml
from pptx import Presentation


class PPTXRewardEvaluator:
    """Evaluates proofreading rewards for PPTX files based on content analysis."""
    
    def __init__(self, supported_languages_path: str = None):
        """
        Initialize the PPTX reward evaluator.
        
        Args:
            supported_languages_path: Path to the supported languages YAML file
        """
        self.euros_per_word = 0.0006  # 0.06 cents per word
        self.language_factors = {}
        
        # Load language difficulty factors
        if supported_languages_path is None:
            supported_languages_path = os.path.join(
                os.path.dirname(os.path.dirname(__file__)), 
                'supported_languages.yml'
            )
        
        self.load_language_factors(supported_languages_path)
    
    def load_language_factors(self, path: str):
        """Load language difficulty factors from YAML file."""
        try:
            with open(path, 'r', encoding='utf-8') as file:
                self.language_factors = yaml.safe_load(file) or {}
        except Exception as e:
            print(f"Warning: Could not load language factors from {path}: {e}")
            # Default to English factor if loading fails
            self.language_factors = {'en': 1.5}
    
    def count_text_boxes_and_words(self, slide) -> Tuple[int, int]:
        """
        Count text boxes and total words in a slide using the same approach as PPTX translation tool.
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            Tuple of (number_of_text_boxes, total_words)
        """
        text_boxes = 0
        total_words = 0
        
        for shape in slide.shapes:
            # Check if shape has a text frame and text (same as PPTX translation tool)
            if not shape.has_text_frame or not shape.text_frame.text.strip():
                continue  # Skip shapes without text
            
            text_boxes += 1
            
            # Count words in the text frame
            text_frame = shape.text_frame
            full_text = text_frame.text
            if full_text.strip():
                # Count words by splitting on whitespace
                words = full_text.split()
                total_words += len(words)
        
        return text_boxes, total_words
    
    def detect_presentation_type(self, presentation) -> str:
        """
        Detect if presentation is image-heavy or video-heavy.
        
        Args:
            presentation: PowerPoint presentation object
            
        Returns:
            'image' or 'video' based on content analysis
        """
        # Simple heuristic: count non-text shapes vs text shapes
        non_text_shapes = 0
        text_shapes = 0
        
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    text_shapes += 1
                else:
                    non_text_shapes += 1
        
        # If more non-text shapes (likely images/media), assume image presentation
        # Otherwise, assume video presentation (more text-based)
        return 'image' if non_text_shapes > text_shapes else 'video'
    
    def calculate_slide_reward(self, text_boxes: int, words: int, 
                             language_factor: float, image_factor: float) -> float:
        """
        Calculate reward for a single slide.
        
        Args:
            text_boxes: Number of text boxes in the slide
            words: Total word count in the slide
            language_factor: Language difficulty multiplier
            image_factor: Image/video presentation factor
            
        Returns:
            Reward amount in euros
        """
        # Delta: 0 if no words, 1 if words exist
        delta = 1 if words > 0 else 0
        
        # Calculate text box fee based on image factor
        text_box_fee = 0.01 * image_factor  # Base fee adjusted by image factor
        
        # Apply formula: reward = delta * (words * diff_lang * euros_per_word + text_box_fee * nbr_text_box)
        reward = delta * language_factor * (
            words * self.euros_per_word + 
            text_box_fee * text_boxes
        )
        
        return reward
    
    def evaluate_pptx(self, file_path: str, language_code: str = 'en', 
                     mode: str = 'auto') -> Dict:
        """
        Evaluate a single PPTX file for proofreading reward.
        
        Args:
            file_path: Path to the PPTX file
            language_code: Language code for difficulty factor
            mode: 'image', 'video', or 'auto' for automatic detection
            
        Returns:
            Dictionary with evaluation results
        """
        try:
            presentation = Presentation(file_path)
            
            # Get language factor
            language_factor = self.language_factors.get(language_code, 1.5)
            
            # Determine image factor
            if mode == 'auto':
                detected_type = self.detect_presentation_type(presentation)
                image_factor = 1.5 if detected_type == 'image' else 1.0
            else:
                image_factor = 1.5 if mode == 'image' else 1.0
            
            # Process each slide
            slide_results = []
            total_reward = 0.0
            total_text_boxes = 0
            total_words = 0
            
            for i, slide in enumerate(presentation.slides, 1):
                text_boxes, words = self.count_text_boxes_and_words(slide)
                slide_reward = self.calculate_slide_reward(
                    text_boxes, words, language_factor, image_factor
                )
                
                slide_results.append({
                    'slide_number': i,
                    'text_boxes': text_boxes,
                    'words': words,
                    'reward': slide_reward
                })
                
                total_reward += slide_reward
                total_text_boxes += text_boxes
                total_words += words
            
            return {
                'file_path': file_path,
                'filename': os.path.basename(file_path),
                'language_code': language_code,
                'language_factor': language_factor,
                'mode': mode if mode != 'auto' else detected_type,
                'image_factor': image_factor,
                'total_slides': len(presentation.slides),
                'total_text_boxes': total_text_boxes,
                'total_words': total_words,
                'total_reward': total_reward,
                'slide_details': slide_results
            }
            
        except Exception as e:
            return {
                'file_path': file_path,
                'filename': os.path.basename(file_path),
                'error': str(e),
                'total_reward': 0.0
            }
    
    def evaluate_multiple_pptx(self, file_paths: List[str], language_code: str = 'en', 
                              mode: str = 'auto') -> List[Dict]:
        """
        Evaluate multiple PPTX files.
        
        Args:
            file_paths: List of PPTX file paths
            language_code: Language code for difficulty factor
            mode: 'image', 'video', or 'auto' for automatic detection
            
        Returns:
            List of evaluation results
        """
        results = []
        for file_path in file_paths:
            result = self.evaluate_pptx(file_path, language_code, mode)
            results.append(result)
        return results
    
    def find_pptx_files(self, directory: str, recursive: bool = True) -> List[str]:
        """
        Find all PPTX files in a directory.
        
        Args:
            directory: Directory to search
            recursive: Whether to search subdirectories
            
        Returns:
            List of PPTX file paths
        """
        pptx_files = []
        path = Path(directory)
        
        if recursive:
            pattern = "**/*.pptx"
        else:
            pattern = "*.pptx"
        
        for file_path in path.glob(pattern):
            if file_path.is_file():
                pptx_files.append(str(file_path))
        
        # Also search for .ppt files
        if recursive:
            pattern = "**/*.ppt"
        else:
            pattern = "*.ppt"
        
        for file_path in path.glob(pattern):
            if file_path.is_file():
                pptx_files.append(str(file_path))
        
        return sorted(pptx_files)
    
    def save_results_to_csv(self, results: List[Dict], output_path: str):
        """
        Save evaluation results to CSV file.
        
        Args:
            results: List of evaluation results
            output_path: Path for the output CSV file
        """
        with open(output_path, 'w', newline='', encoding='utf-8') as csvfile:
            fieldnames = ['filename', 'reward', 'total_slides', 'total_text_boxes', 
                         'total_words', 'language_code', 'mode', 'error']
            writer = csv.DictWriter(csvfile, fieldnames=fieldnames)
            
            writer.writeheader()
            for result in results:
                writer.writerow({
                    'filename': result.get('filename', ''),
                    'reward': f"{result.get('total_reward', 0.0):.4f}",
                    'total_slides': result.get('total_slides', 0),
                    'total_text_boxes': result.get('total_text_boxes', 0),
                    'total_words': result.get('total_words', 0),
                    'language_code': result.get('language_code', ''),
                    'mode': result.get('mode', ''),
                    'error': result.get('error', '')
                })
