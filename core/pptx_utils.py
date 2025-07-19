import logging
import shutil
import tempfile
from pathlib import Path
from typing import Dict, List

from pptx import Presentation

logger = logging.getLogger(__name__)

__all__ = ["split_pptx_to_single_slides"]


def _save_single_slide(src_path: Path, dst_path: Path, slide_index: int) -> None:
    """Save only *slide_index* from *src_path* into a new PPTX at *dst_path*.

    Implementation clones the source file then removes all other slides.
    """
    # Work in a temp copy so we don't destroy the original
    with tempfile.TemporaryDirectory() as tmp:
        tmp_copy = Path(tmp) / "work.pptx"
        shutil.copy(src_path, tmp_copy)
        prs = Presentation(str(tmp_copy))

        # Remove slides except requested index
        sldIdLst = prs.slides._sldIdLst
        sldId_elements = list(sldIdLst)

        for idx, sldId in reversed(list(enumerate(sldId_elements))):
            if idx != slide_index:
                rId = sldId.rId
                prs.part.drop_rel(rId)
                sldIdLst.remove(sldId)
         
        prs.save(dst_path)


def split_pptx_to_single_slides(src_path: Path, out_dir: Path, output_filenames: List[str]) -> List[Path]:
    """Split *src_path* into multiple one-slide presentations.

    Args:
        src_path: translated PPTX containing N slides.
        out_dir: directory to write output files.
        output_filenames: list of filenames to assign per slide, length must equal slide count.

    Returns:
        List of Paths created (same order as slides).
    """
    prs = Presentation(str(src_path))
    slide_count = len(prs.slides)
    if slide_count != len(output_filenames):
        raise ValueError(f"Expecting {slide_count} output filenames, got {len(output_filenames)}")

    out_paths: List[Path] = []
    for idx, fname in enumerate(output_filenames):
        dst = out_dir / fname
        _save_single_slide(src_path, dst, idx)
        out_paths.append(dst)
        logger.info("[PPTX] Saved slide %d -> %s", idx, dst)
    return out_paths 