"""
PPTX to PDF/PNG/WEBP Conversion Core Module

This module provides PowerPoint to PDF, PNG, and WEBP conversion functionality using ConvertAPI.
It can convert PPTX presentations to high-quality PDF documents or individual PNG/WEBP images.

Usage Examples:
    GUI: Integrated with file browsers, format selection, and progress indicators
    API: Used via REST endpoints for document conversion services
    CLI: Command-line presentation conversion for batch processing

Features:
    - ConvertAPI integration for reliable conversion
    - PDF conversion (single document output)
    - PNG conversion (individual slide images)
    - WEBP conversion (PNG to WEBP with optimization)
    - High-quality output preservation
    - Progress callback support for user feedback
    - Comprehensive error handling and validation
    - Batch processing capabilities

Output Formats:
    - PDF: Single document with all slides
    - PNG: Individual image files per slide
    - WEBP: Individual optimized image files per slide
    - Maintains original slide dimensions and quality
    - Preserves fonts and formatting

Conversion Quality:
    - High-resolution output
    - Vector graphics preservation (PDF)
    - WEBP optimization with quality control
    - Consistent rendering across platforms
    - Professional document quality
"""

import logging
import convertapi
from pathlib import Path
from typing import Optional, Callable, List, Tuple
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

class PPTXConverterCore:
    """
    Core PPTX to PDF/PNG/WEBP conversion functionality using ConvertAPI.

    This class provides reliable PowerPoint conversion capabilities with
    support for PDF, PNG, and WEBP output formats. It maintains high quality
    and preserves original presentation formatting.

    Key Features:
        - Reliable conversion via ConvertAPI
        - Multiple output format support (PDF, PNG, WEBP)
        - High-quality output preservation
        - WEBP optimization with quality control
        - Progress tracking with callback support
        - Batch processing capabilities
        - Comprehensive error handling
        - Professional document quality

    Requirements:
        - ConvertAPI key and subscription
        - convertapi Python library
        - PIL/Pillow for WEBP conversion
        - Valid PPTX input files

    Example Usage:
        converter = PPTXConverterCore(api_key="your-convertapi-key")

        # Convert to PDF
        success = converter.convert_pptx_to_pdf(
            input_path=Path("presentation.pptx"),
            output_path=Path("document.pdf")
        )

        # Convert to PNG images
        png_files = converter.convert_pptx_to_png(
            input_path=Path("presentation.pptx"),
            output_dir=Path("images/")
        )

        # Convert to WEBP images
        webp_files = converter.convert_pptx_to_webp(
            input_path=Path("presentation.pptx"),
            output_dir=Path("images/")
        )
    """

    def __init__(self, api_key: str, progress_callback: Optional[Callable[[str], None]] = None):
        """
        Initialize PPTX converter core.

        Args:
            api_key: ConvertAPI key
            progress_callback: Optional callback function for progress updates
        """
        self.api_key = api_key
        self.progress_callback = progress_callback or (lambda x: None)
        self._init_convert_api()

    def _init_convert_api(self):
        """Initialize ConvertAPI."""
        if not self.api_key:
            raise ValueError("ConvertAPI key is required")

        if not isinstance(self.api_key, str) or len(self.api_key.strip()) < 10:
            raise ValueError("ConvertAPI key appears to be invalid (too short or not a string)")

        # ConvertAPI Python SDK uses api_credentials for authentication
        api_key_cleaned = self.api_key.strip()

        # Set the credentials directly (this is what the library actually uses)
        convertapi.api_credentials = api_key_cleaned
        self.progress_callback(f"ConvertAPI initialized with credentials (length: {len(api_key_cleaned)})")
        self.progress_callback("ConvertAPI ready for conversions")

    def convert_pptx_to_pdf(self, input_path: Path, output_path: Path) -> bool:
        """
        Convert PPTX file to PDF.

        Args:
            input_path: Path to input PPTX file
            output_path: Path to output PDF file

        Returns:
            True if successful, False otherwise
        """
        try:
            self.progress_callback(f"Converting PPTX to PDF: {input_path}")

            # Validate input file
            if not self.validate_pptx_file(input_path):
                raise ValueError(f"Invalid PPTX file: {input_path}")

            # Convert using ConvertAPI
            result = convertapi.convert('pdf', {
                'File': str(input_path)
            }, from_format='pptx')

            # Save result
            self.progress_callback(f"Saving PDF to: {output_path}")
            result.file.save(str(output_path))

            self.progress_callback("PPTX to PDF conversion completed successfully")
            return True

        except Exception as e:
            error_msg = f"Failed to convert PPTX to PDF: {e}"
            logger.error(error_msg)
            self.progress_callback(f"Error: {error_msg}")
            return False

    def _calculate_slide_content_bounds(self, slide) -> Optional[Tuple[int, int, int, int]]:
        """
        Calculate the bounding box of all non-background shapes on a slide.
        
        Args:
            slide: python-pptx slide object
            
        Returns:
            Tuple of (left, top, right, bottom) in EMU units or None if no shapes
        """
        if not slide.shapes:
            return None
            
        min_left = float('inf')
        min_top = float('inf')
        max_right = float('-inf')
        max_bottom = float('-inf')
        
        has_content = False
        
        for shape in slide.shapes:
            # Skip background shapes and placeholders without content
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                if not hasattr(shape, 'text') or not shape.text.strip():
                    continue
            
            # Get shape bounds
            try:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                
                if left is not None and top is not None and width and height:
                    has_content = True
                    min_left = min(min_left, left)
                    min_top = min(min_top, top)
                    max_right = max(max_right, left + width)
                    max_bottom = max(max_bottom, top + height)
                    
                    # Handle group shapes recursively
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        for child_shape in shape.shapes:
                            child_left = child_shape.left
                            child_top = child_shape.top
                            child_width = child_shape.width
                            child_height = child_shape.height
                            
                            if child_left is not None and child_top is not None and child_width and child_height:
                                min_left = min(min_left, child_left)
                                min_top = min(min_top, child_top)
                                max_right = max(max_right, child_left + child_width)
                                max_bottom = max(max_bottom, child_top + child_height)
            except Exception as e:
                logger.debug(f"Skipping shape due to error: {e}")
                continue
        
        if not has_content:
            return None
            
        return (min_left, min_top, max_right, max_bottom)
    
    def _crop_png_to_content(self, png_path: Path, bounds: Optional[Tuple[int, int, int, int]] = None, 
                            slide_width: int = None, slide_height: int = None) -> bool:
        """
        Crop a PNG image to its content bounds or specified bounds.
        
        Args:
            png_path: Path to PNG file to crop
            bounds: Optional tuple of (left, top, right, bottom) in EMU units
            slide_width: Slide width in EMU units for scaling
            slide_height: Slide height in EMU units for scaling
            
        Returns:
            True if cropping was successful
        """
        try:
            img = Image.open(png_path)
            
            if img.mode != 'RGBA':
                img = img.convert('RGBA')
            
            if bounds and slide_width and slide_height:
                # Convert EMU bounds to pixel coordinates
                img_width, img_height = img.size
                
                # Calculate scale factors
                scale_x = img_width / slide_width
                scale_y = img_height / slide_height
                
                # Convert bounds to pixels
                left = int(bounds[0] * scale_x)
                top = int(bounds[1] * scale_y)
                right = int(bounds[2] * scale_x)
                bottom = int(bounds[3] * scale_y)
                
                # Ensure bounds are within image
                left = max(0, left)
                top = max(0, top)
                right = min(img_width, right)
                bottom = min(img_height, bottom)
                
                # Crop the image
                if left < right and top < bottom:
                    cropped = img.crop((left, top, right, bottom))
                else:
                    # If bounds are invalid, try auto-crop
                    bbox = img.getbbox()
                    if bbox:
                        cropped = img.crop(bbox)
                    else:
                        return False
            else:
                # Auto-crop to content using getbbox
                bbox = img.getbbox()
                if bbox:
                    cropped = img.crop(bbox)
                else:
                    # Image is fully transparent or empty
                    return False
            
            # Save the cropped image
            cropped.save(png_path, 'PNG')
            return True
            
        except Exception as e:
            logger.error(f"Failed to crop PNG {png_path}: {e}")
            return False

    def convert_pptx_to_png(self, input_path: Path, output_dir: Path, group_elements: bool = False) -> List[str]:
        """
        Convert PPTX file to PNG images (one per slide).

        Args:
            input_path: Path to input PPTX file
            output_dir: Directory to save PNG files
            group_elements: If True, crop PNG to content bounds (default: False for backward compatibility)

        Returns:
            List of paths to generated PNG files
        """
        try:
            self.progress_callback(f"Converting PPTX to PNG: {input_path}")

            # Validate input file
            if not self.validate_pptx_file(input_path):
                raise ValueError(f"Invalid PPTX file: {input_path}")

            # Ensure output directory exists
            output_dir.mkdir(parents=True, exist_ok=True)

            # Get base name safely
            base_name = input_path.stem if input_path.stem else "slide"
            self.progress_callback(f"Using base name: {base_name}")

            # Convert using ConvertAPI
            self.progress_callback("Calling ConvertAPI...")
            result = convertapi.convert('png', {
                'File': str(input_path)
            }, from_format='pptx')

            # Validate result
            if not result:
                raise RuntimeError("ConvertAPI returned empty result")

            # Save all PNG files
            png_files = []
            self.progress_callback("Processing conversion results...")

            if hasattr(result, 'files') and result.files:
                # Multiple files (slides)
                self.progress_callback(f"Processing {len(result.files)} slide files")
                for i, file_result in enumerate(result.files, 1):
                    if not file_result:
                        self.progress_callback(f"Warning: Slide {i} file result is None, skipping")
                        continue

                    png_path = output_dir / f"{base_name}_slide_{i:02d}.png"
                    self.progress_callback(f"Saving slide {i} to: {png_path}")

                    try:
                        file_result.save(str(png_path))
                        if png_path.exists() and png_path.stat().st_size > 0:
                            png_files.append(str(png_path))
                            self.progress_callback(f"Successfully saved slide {i} as PNG")
                        else:
                            self.progress_callback(f"Warning: Slide {i} was not saved properly")
                    except Exception as save_error:
                        self.progress_callback(f"Error saving slide {i}: {save_error}")
                        continue
            else:
                # Single file
                self.progress_callback("Processing single file result")
                if hasattr(result, 'file') and result.file:
                    png_path = output_dir / f"{base_name}.png"
                    self.progress_callback(f"Saving single file to: {png_path}")

                    try:
                        result.file.save(str(png_path))
                        if png_path.exists() and png_path.stat().st_size > 0:
                            png_files.append(str(png_path))
                            self.progress_callback("Successfully saved PNG file")
                        else:
                            self.progress_callback("Warning: Single file was not saved properly")
                    except Exception as save_error:
                        self.progress_callback(f"Error saving single file: {save_error}")
                else:
                    raise RuntimeError("ConvertAPI result has no files or file attribute")

            if not png_files:
                raise RuntimeError("No PNG files were generated successfully")

            # Apply element grouping/cropping if requested
            if group_elements:
                self.progress_callback("Applying element grouping to PNG files...")
                
                try:
                    # Load the PPTX to get slide dimensions and element bounds
                    prs = Presentation(str(input_path))
                    
                    for i, png_file in enumerate(png_files):
                        if i < len(prs.slides):
                            slide = prs.slides[i]
                            
                            # Get slide dimensions
                            slide_width = prs.slide_width
                            slide_height = prs.slide_height
                            
                            # Calculate content bounds for this slide
                            bounds = self._calculate_slide_content_bounds(slide)
                            
                            # Crop the PNG to content bounds
                            png_path = Path(png_file)
                            if png_path.exists():
                                if bounds:
                                    # Crop to calculated bounds
                                    success = self._crop_png_to_content(png_path, bounds, slide_width, slide_height)
                                    if success:
                                        self.progress_callback(f"Cropped slide {i+1} to element bounds")
                                    else:
                                        # Fallback to auto-crop
                                        success = self._crop_png_to_content(png_path)
                                        if success:
                                            self.progress_callback(f"Auto-cropped slide {i+1} to content")
                                else:
                                    # No bounds calculated, try auto-crop
                                    success = self._crop_png_to_content(png_path)
                                    if success:
                                        self.progress_callback(f"Auto-cropped slide {i+1} to content")
                
                except Exception as e:
                    # Log the error but don't fail - images are already converted
                    logger.warning(f"Could not apply element grouping: {e}")
                    self.progress_callback(f"Warning: Element grouping failed, using full slides: {e}")

            self.progress_callback(f"PPTX to PNG conversion completed successfully. Generated {len(png_files)} files")
            return png_files

        except Exception as e:
            error_msg = f"Failed to convert PPTX to PNG: {e}"
            logger.error(error_msg, exc_info=True)
            self.progress_callback(f"Error: {error_msg}")
            return []

    def validate_pptx_file(self, file_path: Path) -> bool:
        """Validate that the file is a valid PPTX file."""
        if not file_path.exists():
            return False

        if file_path.suffix.lower() != '.pptx':
            return False

        # Basic file size check (PPTX files should be larger than empty ZIP)
        return file_path.stat().st_size > 1000

    def convert_pptx_to_webp(self, input_path: Path, output_dir: Path, quality: int = 85, group_elements: bool = False) -> List[str]:
        """
        Convert PPTX file to WEBP images (one per slide).

        Args:
            input_path: Path to input PPTX file
            output_dir: Directory to save WEBP files
            quality: WEBP quality (0-100, default 85)
            group_elements: If True, crop images to content bounds (default: False)

        Returns:
            List of paths to generated WEBP files
        """
        try:
            self.progress_callback(f"Converting PPTX to WEBP: {input_path}")

            # Validate input file
            if not self.validate_pptx_file(input_path):
                raise ValueError(f"Invalid PPTX file: {input_path}")

            # Ensure output directory exists
            output_dir.mkdir(parents=True, exist_ok=True)

            # First convert to PNG using existing method (with group_elements if requested)
            self.progress_callback("Converting to PNG first...")
            png_files = self.convert_pptx_to_png(input_path, output_dir, group_elements)

            if not png_files:
                raise RuntimeError("Failed to convert PPTX to PNG")

            # Convert each PNG to WEBP
            webp_files = []
            base_name = input_path.stem

            self.progress_callback("Converting PNG files to WEBP...")

            for i, png_path in enumerate(png_files, 1):
                png_file = Path(png_path)
                webp_name = f"{base_name}_slide_{i:02d}.webp"
                webp_path = output_dir / webp_name

                # Convert PNG to WEBP using PIL
                self.progress_callback(f"Converting slide {i+1} to WEBP...")
                with Image.open(png_file) as img:
                    img.save(webp_path, 'WEBP', quality=quality, method=6)

                # Remove the temporary PNG file
                png_file.unlink()

                webp_files.append(str(webp_path))
                self.progress_callback(f"Saved slide {i+1} as WEBP")

            self.progress_callback(f"PPTX to WEBP conversion completed successfully. Generated {len(webp_files)} files")
            return webp_files

        except Exception as e:
            error_msg = f"Failed to convert PPTX to WEBP: {e}"
            logger.error(error_msg)
            self.progress_callback(f"Error: {error_msg}")
            return []

    def get_supported_formats(self) -> List[str]:
        """Get list of supported output formats."""
        return ['pdf', 'png', 'webp']
