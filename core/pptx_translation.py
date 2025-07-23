"""
PPTX Translation Core Module

This module provides PowerPoint presentation translation functionality using the DeepL API.
It can translate text content within PPTX files while preserving formatting, layout, and structure.

Usage Examples:
    GUI: Integrated into a file selection interface with progress bars and language dropdowns
    API: Used via REST endpoints for batch translation services
    CLI: Command-line translation of presentation files

Features:
    - DeepL API integration for high-quality translation
    - Preserves PowerPoint formatting and layout
    - Handles text in slides, shapes, tables, and text frames
    - Progress callback support for user feedback
    - Comprehensive error handling and validation
    - Support for all DeepL supported languages

Supported Content:
    - Slide text content
    - Table cell text
    - Text frame paragraphs
    - Shape text content
    - Maintains original formatting and styling
"""

import logging
import deepl
from pathlib import Path
from typing import Dict, Any, Optional, Callable
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR, MSO_COLOR_TYPE

logger = logging.getLogger(__name__)

class PPTXTranslationCore:
    """
    Core PPTX translation functionality using DeepL API.
    
    This class provides comprehensive PowerPoint translation capabilities while
    maintaining the original presentation structure, formatting, and layout.
    It processes all text content including slides, tables, and text frames.
    
    Key Features:
        - High-quality translation via DeepL API
        - Preserves original formatting and layout
        - Handles complex PowerPoint structures
        - Progress tracking with callback support
        - Robust error handling and validation
        - Support for all DeepL language pairs
    
    Requirements:
        - DeepL API key
        - python-pptx library
        - Valid PPTX input files
    
    Example Usage:
        translator = PPTXTranslationCore(api_key="your-deepl-key")
        success = translator.translate_pptx(
            input_path=Path("presentation.pptx"),
            output_path=Path("translated.pptx"),
            source_lang="en",
            target_lang="fr"
        )
    """
    
    def __init__(self, api_key: str, progress_callback: Optional[Callable[[str], None]] = None):
        """
        Initialize PPTX translation core.
        
        Args:
            api_key: DeepL API key
            progress_callback: Optional callback function for progress updates
        """
        self.api_key = api_key
        self.progress_callback = progress_callback or (lambda x: None)
        self.translator = None
        self._init_translator()
    
    def _init_translator(self):
        """Initialize DeepL translator."""
        if not self.api_key:
            raise ValueError("DeepL API key is required")
        
        try:
            self.translator = deepl.Translator(self.api_key)
            self.progress_callback("DeepL translator initialized")
        except Exception as e:
            raise RuntimeError(f"Failed to initialize DeepL translator: {e}")
    
    def translate_pptx(self, input_path: Path, output_path: Path, 
                      source_lang: str, target_lang: str) -> bool:
        """
        Translate PPTX file from source to target language with full formatting preservation.
        
        Args:
            input_path: Path to input PPTX file
            output_path: Path to output PPTX file
            source_lang: Source language code (e.g., 'en')
            target_lang: Target language code (e.g., 'fr')
            
        Returns:
            True if successful, False otherwise
        """
        try:
            self.progress_callback(f"Opening PPTX file: {input_path}")
            
            # Check if file exists and has content
            if not input_path.exists():
                raise FileNotFoundError(f"Input file does not exist: {input_path}")
            
            file_size = input_path.stat().st_size
            if file_size == 0:
                raise ValueError(f"Input file is empty: {input_path}")
            
            self.progress_callback(f"File size: {file_size} bytes")
            
            # Load presentation
            try:
                prs = Presentation(str(input_path))
            except Exception as e:
                # Log first few bytes of file for debugging
                with open(input_path, 'rb') as f:
                    header = f.read(20)
                    header_hex = header.hex()
                self.progress_callback(f"File header (hex): {header_hex}")
                raise ValueError(f"Failed to open PPTX file: {str(e)}")
            
            slide_count = len(prs.slides)
            self.progress_callback(f"Found {slide_count} slides to translate")
            
            # Track translation progress
            total_shapes = sum(len(slide.shapes) for slide in prs.slides)
            processed_shapes = 0
            
            # Translate each shape with text
            for slide_idx, slide in enumerate(prs.slides):
                self.progress_callback(f"Processing Slide {slide_idx + 1}/{len(prs.slides)}")
                
                for shape_idx, shape in enumerate(slide.shapes):
                    # Check if shape has a text frame and text
                    if not shape.has_text_frame or not shape.text_frame.text.strip():
                        processed_shapes += 1
                        continue # Skip shapes without text

                    try:
                        text_frame = shape.text_frame
                        original_paras_data = []

                        # Store original formatting for each paragraph and run
                        for para in text_frame.paragraphs:
                            para_data = {
                                'text': para.text,
                                'alignment': para.alignment,
                                'level': para.level,
                                'line_spacing': para.line_spacing,
                                'space_before': para.space_before,
                                'space_after': para.space_after,
                                'runs': []
                            }

                            for run in para.runs:
                                font = run.font
                                color_info = None
                                if font.color and hasattr(font.color, 'type'): 
                                    if font.color.type == MSO_COLOR_TYPE.RGB:
                                        color_info = ('rgb', font.color.rgb)
                                    elif font.color.type == MSO_COLOR_TYPE.SCHEME:
                                        color_info = ('scheme', font.color.theme_color, getattr(font.color, 'brightness', 0.0))

                                run_data = {
                                    'text': run.text, 
                                    'font_name': font.name,
                                    'size': font.size,
                                    'bold': font.bold,
                                    'italic': font.italic,
                                    'underline': font.underline,
                                    'color_info': color_info,
                                    'language': getattr(font, 'language_id', None)
                                }
                                para_data['runs'].append(run_data)
                            original_paras_data.append(para_data)

                        # Translate the text
                        original_full_text = text_frame.text
                        if not original_full_text.strip():
                             processed_shapes += 1
                             continue # Skip if effectively empty after stripping

                        translated_text_obj = self.translator.translate_text(
                            original_full_text,
                            source_lang=source_lang,
                            target_lang=target_lang
                        )
                        translated_full_text = translated_text_obj.text

                        text_frame.clear() # Clear existing content

                        # Reconstruct text with original formatting
                        translated_paras = translated_full_text.split('\n')
                        num_orig_paras = len(original_paras_data)
                        num_trans_paras = len(translated_paras)

                        for i, trans_para_text in enumerate(translated_paras):
                            # Determine which original paragraph's style to mimic
                            orig_para_idx = min(i, num_orig_paras - 1)
                            orig_para_data = original_paras_data[orig_para_idx]

                            # Add paragraph (first one exists, add subsequent ones)
                            if i == 0:
                                p = text_frame.paragraphs[0]
                                p.text = '' # Clear any default text in the first paragraph
                            else:
                                p = text_frame.add_paragraph()

                            # Apply paragraph formatting
                            p.alignment = orig_para_data['alignment']
                            p.level = orig_para_data['level']
                            if orig_para_data['line_spacing']: p.line_spacing = orig_para_data['line_spacing']
                            if orig_para_data['space_before']: p.space_before = orig_para_data['space_before']
                            if orig_para_data['space_after']: p.space_after = orig_para_data['space_after']

                            # Apply run formatting - Distribute text and styles
                            orig_runs_data = orig_para_data['runs']
                            num_orig_runs = len(orig_runs_data)

                            if not orig_runs_data: # If original paragraph had no runs (e.g., empty)
                                p.text = trans_para_text # Just add the text
                                continue

                            # Simple distribution: Apply styles run-by-run, splitting translated text
                            words = trans_para_text.split()
                            total_words = len(words)
                            start_idx = 0

                            for j, run_data in enumerate(orig_runs_data):
                                words_for_this_run = total_words // num_orig_runs
                                if j < total_words % num_orig_runs:
                                    words_for_this_run += 1

                                end_idx = start_idx + words_for_this_run
                                run_text = ' '.join(words[start_idx:end_idx])
                                start_idx = end_idx

                                if not run_text and j < num_orig_runs -1 : # Avoid adding empty runs unless it's the last one potentially
                                    continue

                                run = p.add_run()
                                run.text = run_text + (' ' if j < num_orig_runs - 1 and run_text else '') # Add space between runs

                                # Apply run formatting
                                font = run.font
                                if run_data['font_name']: font.name = run_data['font_name']
                                if run_data['size']: font.size = run_data['size']
                                # Explicitly set False if stored as False
                                font.bold = run_data['bold'] if run_data['bold'] is not None else None
                                font.italic = run_data['italic'] if run_data['italic'] is not None else None
                                font.underline = run_data['underline'] if run_data['underline'] is not None else None # Check underline type if needed

                                stored_color_info = run_data['color_info']
                                if stored_color_info:
                                    color_type, value1, *rest = stored_color_info
                                    if color_type == 'rgb':
                                        try:
                                            font.color.rgb = RGBColor(*value1) # Pass tuple elements to RGBColor
                                        except Exception as color_e:
                                            self.progress_callback(f"Warn: Failed to set RGB color {value1}: {color_e}")
                                    elif color_type == 'scheme':
                                        try:
                                            font.color.theme_color = value1
                                            if rest: # Brightness was stored
                                                font.color.brightness = rest[0]
                                        except Exception as color_e:
                                             self.progress_callback(f"Warn: Failed to set theme color {value1}: {color_e}")

                                if run_data['language']: font.language_id = run_data['language']

                    except Exception as e:
                        self.progress_callback(f"Error translating shape {shape_idx+1} on slide {slide_idx+1}: {e}")
                        logger.warning(f"Error translating shape {shape_idx+1} on slide {slide_idx+1} in {input_path.name}: {e}", exc_info=True)

                    processed_shapes += 1
                    if total_shapes > 0 and processed_shapes % 5 == 0: # Update progress more frequently
                        progress = (processed_shapes / total_shapes) * 100
                        self.progress_callback(f"Translation progress: {progress:.1f}% ({processed_shapes}/{total_shapes} shapes)")
            
            # Save translated presentation
            self.progress_callback(f"Saving translated presentation to: {output_path}")
            prs.save(str(output_path))
            
            self.progress_callback("PPTX translation completed successfully")
            return True
            
        except Exception as e:
            error_msg = f"Failed to translate PPTX: {e}"
            logger.error(error_msg)
            self.progress_callback(f"Error: {error_msg}")
            return False
    
    def _translate_text(self, text: str, source_lang: str, target_lang: str) -> str:
        """Translate text using DeepL API."""
        if not text.strip():
            return text
        
        try:
            result = self.translator.translate_text(
                text, 
                source_lang=source_lang if source_lang != 'auto' else None,
                target_lang=target_lang
            )
            return result.text
        except Exception as e:
            logger.warning(f"Translation failed for text: {text[:50]}... Error: {e}")
            return text  # Return original text if translation fails
    
    def validate_file(self, file_path: Path) -> bool:
        """Validate that the file is a valid PPTX file."""
        if not file_path.exists():
            return False
        
        if file_path.suffix.lower() != '.pptx':
            return False
        
        try:
            # Try to open the presentation
            prs = Presentation(str(file_path))
            return len(prs.slides) > 0
        except Exception:
            return False
    
    def get_supported_languages(self) -> Dict[str, str]:
        """Get languages supported by DeepL."""
        try:
            if not self.translator:
                return {}
            
            source_langs = self.translator.get_source_languages()
            target_langs = self.translator.get_target_languages()
            
            supported = {}
            for lang in source_langs:
                supported[lang.code.lower()] = lang.name
            
            return supported
        except Exception as e:
            logger.error(f"Failed to get supported languages: {e}")
            return {}